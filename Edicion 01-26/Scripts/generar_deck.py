#!/usr/bin/env python3
"""
generar_deck.py — IA en la Oficina · Edición 01/26

Genera el PPTX del formador a partir del brief de la sesión.

REGLA DE ORO INNEGOCIABLE: todo el contenido textual del PPTX viene del brief.
Este script no inventa pedagogía, ejemplos, prompts, datos, citas, glosario ni
notas del orador. Si el brief no lo dice, el PPTX no lo lleva.

Uso:
    python3 generar_deck.py D05
    python3 generar_deck.py D01 D02 D03
    python3 generar_deck.py --all
    python3 generar_deck.py D05 --validate
    python3 generar_deck.py D05 --dump-ir D05_ir.json
    python3 generar_deck.py D05 --out /tmp/foo.pptx

Exit codes:
    0  éxito
    1  warnings (deck generado igualmente)
    2  error fatal de validación
    3  error de I/O
"""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
from dataclasses import dataclass, field, asdict
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any

# ─── config ──────────────────────────────────────────────────────────────────

SCRIPT_DIR = Path(__file__).resolve().parent
ROOT_DIR = SCRIPT_DIR.parent
BRIEFS_DIR = ROOT_DIR / "Briefs"
DECKS_DIR = ROOT_DIR / "Decks"
LOGOS_DIR = ROOT_DIR / "Logos"
BANK_PATH = ROOT_DIR / "Recursos" / "Banco_de_Ejercicios.md"

LOGO_PATH = LOGOS_DIR / "00_FULL COLOR (1).png"

# Paleta AllWomen (RGB)
COLOR_AZUL = (0x2B, 0x5C, 0xE6)
COLOR_NARANJA = (0xE8, 0x60, 0x1A)
COLOR_CREMA = (0xF5, 0xEF, 0xE6)
COLOR_OSCURO = (0x1A, 0x1A, 0x2E)
COLOR_GRIS = (0x66, 0x66, 0x66)
COLOR_BLANCO = (0xFF, 0xFF, 0xFF)

# Dimensiones del deck (pulgadas)
SLIDE_W_IN = 17.78
SLIDE_H_IN = 10.0

FOOTER_TEXT = (
    "IA en la Oficina · Edición 01/26 · "
    "Fundación ONCE · Inserta Empleo · Talento Digital"
)
SECCIONES_OBLIGATORIAS = {
    2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 15, 16, 17, 19, 20, 21, 22
}
DIAS_CON_HILO_PROYECTO = {"D01", "D09", "D16", "D17"}

MAX_BULLETS_POR_SLIDE = 5
MAX_PALABRAS_POR_BULLET = 15
MAX_SLIDES = 80

ICONS_BY_LAYOUT = {
    "portada": "",
    "agenda": "🗓",
    "recap": "🔁",
    "objetivos": "🎯",
    "concepto": "💡",
    "datos": "📊",
    "ejemplo": "🔍",
    "errores": "⚠️",
    "glosario": "📖",
    "practica": "🛠",
    "ficha_ejercicio": "✏️",
    "parejas": "👥",
    "puesta_comun": "🗣",
    "mensaje_clave": "✨",
    "cierre_hoy": "✅",
    "cierre_manana": "🎓",
    "proyecto_final": "🎓",
}

# ─── models ──────────────────────────────────────────────────────────────────

@dataclass
class Concept:
    title: str
    paragraphs: List[str] = field(default_factory=list)
    bullets_hint: List[str] = field(default_factory=list)  # **Puntos clave:**

@dataclass
class Example:
    title: str
    situacion: str = ""
    prompt: str = ""
    resultado: str = ""
    iteracion: str = ""

@dataclass
class GlossaryTerm:
    term: str
    definition: str

@dataclass
class Activity:
    title: str
    modalidad: str = ""
    duracion: str = ""
    descripcion: str = ""
    objetivo: str = ""
    material: str = ""

@dataclass
class QuizItem:
    number: int
    question: str
    options: List[str] = field(default_factory=list)  # con ✓ en correcta si hay
    answer: str = ""  # respuesta abierta o explicación

@dataclass
class FAQItem:
    question: str
    answer: str

@dataclass
class Tool:
    name: str
    description: str = ""

@dataclass
class Brief:
    code: str
    path: Path
    title: str = ""
    fecha: str = ""
    contexto: str = ""
    recap_anterior: str = ""
    resumen: str = ""
    estructura_tabla: str = ""
    objetivos: List[str] = field(default_factory=list)
    conceptos: List[Concept] = field(default_factory=list)
    datos: List[str] = field(default_factory=list)
    ejemplos: List[Example] = field(default_factory=list)
    glosario: List[GlossaryTerm] = field(default_factory=list)
    errores: List[str] = field(default_factory=list)
    actividades: List[Activity] = field(default_factory=list)
    parejas: str = ""
    ejercicios: List[str] = field(default_factory=list)
    cuestionario: List[QuizItem] = field(default_factory=list)
    faq: List[FAQItem] = field(default_factory=list)
    herramientas: List[Tool] = field(default_factory=list)
    mini_entregable: str = ""
    hilo_proyecto: str = ""
    mensaje_clave: str = ""
    ideas: List[str] = field(default_factory=list)
    conexion: str = ""
    sections_found: set = field(default_factory=set)

@dataclass
class Slide:
    layout: str  # portada / agenda / recap / objetivos / concepto / ...
    title: str = ""
    subtitle: str = ""
    bullets: List[str] = field(default_factory=list)
    body_text: str = ""
    notes: str = ""
    extras: Dict[str, Any] = field(default_factory=dict)

@dataclass
class ValidationReport:
    brief_code: str
    errors: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)

    def has_errors(self) -> bool:
        return len(self.errors) > 0

    def print_summary(self) -> None:
        if not self.errors and not self.warnings:
            print(f"  ✅ {self.brief_code}: sin incidencias")
            return
        if self.errors:
            print(f"  ❌ {self.brief_code}: {len(self.errors)} error(es):")
            for e in self.errors:
                print(f"     - {e}")
        if self.warnings:
            print(f"  ⚠️  {self.brief_code}: {len(self.warnings)} warning(s):")
            for w in self.warnings:
                print(f"     - {w}")

# ─── parser ──────────────────────────────────────────────────────────────────

RE_SECTION_HEADER = re.compile(r"^##\s+(\d{1,2})\.\s+(.+?)\s*$")
RE_SUBSECTION = re.compile(r"^###\s+(.+?)\s*$")
RE_BOLD_TERM = re.compile(r"^\*\*([^*]+?):\*\*\s*(.*)$")
RE_PUNTOS_CLAVE_HEADER = re.compile(r"^>\s*\*\*Puntos clave:\*\*\s*$", re.IGNORECASE)
RE_BLOCKQUOTE_BULLET = re.compile(r"^>\s*-\s+(.+)$")
RE_BULLET_LINE = re.compile(r"^-\s+(.+)$")
RE_QUIZ_QUESTION = re.compile(r"^\*\*(\d+)\.\s+(.+?)\*\*\s*$")
RE_FAQ_QUESTION = re.compile(r"^\*\*(¿.+?\?.*)\*\*\s*$")
RE_EXAMPLE_FIELD = re.compile(
    r"^\*\*(Situación|Paso\s*\d+\s*[—-]\s*\w+|Prompt[^:*]*|Resultado[^:*]*|Cómo iterarlo)[^:*]*:\*\*\s*(.*)$"
)
RE_EXERCISE_CODE = re.compile(r"\bEJ-D(\d{2})-(\d+)\b")
RE_ACTIVITY_HEADER = re.compile(
    r"^###\s+Actividad\s+\d+\s*[—-]\s*(.+?)\s*\[([^·]+)·\s*([^\]]+)\]\s*$"
)
RE_FECHA = re.compile(r"^\*\*Fecha[:\*]\s*\**(.+?)\s*$")

def _strip_md(s: str) -> str:
    s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)
    s = re.sub(r"\*([^*]+)\*", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    return s.strip()

def parse_brief(path: Path) -> Brief:
    code = path.stem.split("_")[0]  # "D05_Brief_NotebookLM" → "D05"
    text = path.read_text(encoding="utf-8")
    lines = text.splitlines()
    brief = Brief(code=code, path=path)

    # Título de la sesión: primera línea con "## Módulo" o segunda
    for line in lines[:5]:
        if line.startswith("## Módulo") or line.startswith("## Sesión"):
            brief.title = _strip_md(line.lstrip("# ").strip())
            break

    # Fecha
    for line in lines[:15]:
        m = RE_FECHA.match(line)
        if m:
            brief.fecha = _strip_md(m.group(1))
            break

    # Trocear por secciones de nivel 1 (## N. NOMBRE)
    sections: Dict[int, Tuple[str, List[str]]] = {}
    current_num: Optional[int] = None
    current_name: str = ""
    current_lines: List[str] = []
    for line in lines:
        m = RE_SECTION_HEADER.match(line)
        if m:
            if current_num is not None:
                sections[current_num] = (current_name, current_lines)
            current_num = int(m.group(1))
            current_name = m.group(2).strip()
            current_lines = []
            brief.sections_found.add(current_num)
        else:
            if current_num is not None:
                current_lines.append(line)
    if current_num is not None:
        sections[current_num] = (current_name, current_lines)

    if 2 in sections:
        brief.contexto, brief.recap_anterior = _parse_contexto(sections[2][1])
    if 3 in sections:
        brief.resumen = _join_paragraphs(sections[3][1])
    if 4 in sections:
        brief.estructura_tabla = _join_paragraphs(sections[4][1])
    if 5 in sections:
        brief.objetivos = _parse_numbered_list(sections[5][1])
    if 6 in sections:
        brief.conceptos = _parse_conceptos(sections[6][1])
    if 7 in sections:
        brief.datos = _parse_bullets(sections[7][1])
    if 8 in sections:
        brief.ejemplos = _parse_ejemplos(sections[8][1])
    if 9 in sections:
        brief.glosario = _parse_glosario(sections[9][1])
    if 10 in sections:
        brief.errores = _parse_errores(sections[10][1])
    if 11 in sections:
        brief.actividades = _parse_actividades(sections[11][1])
    if 12 in sections:
        brief.parejas = _join_paragraphs(sections[12][1])
    if 13 in sections:
        brief.ejercicios = _parse_bullets(sections[13][1])
    if 14 in sections:
        brief.cuestionario = _parse_cuestionario(sections[14][1])
    if 15 in sections:
        brief.faq = _parse_faq(sections[15][1])
    if 16 in sections:
        brief.herramientas = _parse_herramientas(sections[16][1])
    if 17 in sections:
        brief.mini_entregable = _join_paragraphs(sections[17][1])
    if 18 in sections:
        brief.hilo_proyecto = _join_paragraphs(sections[18][1])
    if 19 in sections:
        brief.mensaje_clave = _parse_mensaje_clave(sections[19][1])
    if 20 in sections:
        brief.ideas = _parse_bullets(sections[20][1])
    if 21 in sections:
        brief.conexion = _join_paragraphs(sections[21][1])

    return brief

def _join_paragraphs(lines: List[str]) -> str:
    out: List[str] = []
    for ln in lines:
        if ln.strip().startswith("---"):
            continue
        out.append(ln)
    return "\n".join(out).strip()

def _parse_contexto(lines: List[str]) -> Tuple[str, str]:
    text = _join_paragraphs(lines)
    recap = ""
    for ln in lines:
        m = re.match(r"^\*\*Sesión anterior[^:]*:\*\*\s*(.+)$", ln)
        if m:
            recap = _strip_md(m.group(1))
            break
    return text, recap

def _parse_numbered_list(lines: List[str]) -> List[str]:
    out: List[str] = []
    for ln in lines:
        m = re.match(r"^\d+\.\s+(.+)$", ln.strip())
        if m:
            out.append(_strip_md(m.group(1)))
    return out

def _parse_bullets(lines: List[str]) -> List[str]:
    out: List[str] = []
    for ln in lines:
        m = RE_BULLET_LINE.match(ln)
        if m:
            out.append(_strip_md(m.group(1)))
    return out

def _parse_conceptos(lines: List[str]) -> List[Concept]:
    conceptos: List[Concept] = []
    current: Optional[Concept] = None
    in_puntos = False
    for ln in lines:
        m = RE_SUBSECTION.match(ln)
        if m:
            if current is not None:
                conceptos.append(current)
            current = Concept(title=_strip_md(m.group(1)))
            in_puntos = False
            continue
        if current is None:
            continue
        if RE_PUNTOS_CLAVE_HEADER.match(ln):
            in_puntos = True
            continue
        if in_puntos:
            mb = RE_BLOCKQUOTE_BULLET.match(ln)
            if mb:
                current.bullets_hint.append(_strip_md(mb.group(1)))
                continue
            # blockquote vacío "> " sin "-" — fin del bloque
            if not ln.strip().startswith(">"):
                in_puntos = False
        if not in_puntos and ln.strip() and not ln.strip().startswith(">"):
            if not ln.strip().startswith("---"):
                current.paragraphs.append(ln.strip())
    if current is not None:
        conceptos.append(current)
    # Compactar párrafos consecutivos en bloques con separador en blanco
    for c in conceptos:
        merged: List[str] = []
        buffer: List[str] = []
        for p in c.paragraphs:
            if p == "":
                if buffer:
                    merged.append(" ".join(buffer))
                    buffer = []
            else:
                buffer.append(p)
        if buffer:
            merged.append(" ".join(buffer))
        c.paragraphs = merged
    return conceptos

def _parse_ejemplos(lines: List[str]) -> List[Example]:
    ejemplos: List[Example] = []
    current: Optional[Example] = None
    current_field: Optional[str] = None
    buffer: List[str] = []

    def flush_field():
        nonlocal buffer, current_field, current
        if current is None or current_field is None:
            buffer = []
            return
        value = " ".join(b.strip() for b in buffer if b.strip()).strip()
        if current_field.startswith("situacion"):
            current.situacion = (current.situacion + " " + value).strip()
        elif current_field.startswith("prompt"):
            current.prompt = (current.prompt + " " + value).strip()
        elif current_field.startswith("resultado"):
            current.resultado = (current.resultado + " " + value).strip()
        elif current_field.startswith("iteracion"):
            current.iteracion = (current.iteracion + " " + value).strip()
        buffer = []

    for ln in lines:
        m = RE_SUBSECTION.match(ln)
        if m:
            flush_field()
            if current is not None:
                ejemplos.append(current)
            current = Example(title=_strip_md(m.group(1)))
            current_field = None
            continue
        if current is None:
            continue
        mf = RE_BOLD_TERM.match(ln)
        if mf:
            flush_field()
            label = mf.group(1).lower()
            content = mf.group(2)
            if label.startswith("situación") or label.startswith("situacion"):
                current_field = "situacion"
            elif label.startswith("prompt"):
                current_field = "prompt"
            elif label.startswith("resultado"):
                current_field = "resultado"
            elif "iterarlo" in label or label.startswith("cómo iterarlo"):
                current_field = "iteracion"
            else:
                current_field = None
            if content:
                buffer.append(_strip_md(content))
            continue
        if current_field is not None and ln.strip() and not ln.strip().startswith("---"):
            # Limpiar blockquotes "> ..."
            cleaned = ln.strip().lstrip(">").strip()
            buffer.append(_strip_md(cleaned))
    flush_field()
    if current is not None:
        ejemplos.append(current)
    return ejemplos

def _parse_glosario(lines: List[str]) -> List[GlossaryTerm]:
    out: List[GlossaryTerm] = []
    for ln in lines:
        m = RE_BOLD_TERM.match(ln)
        if m:
            term = m.group(1).strip()
            definition = _strip_md(m.group(2))
            if term and definition:
                out.append(GlossaryTerm(term=term, definition=definition))
    return out

def _parse_errores(lines: List[str]) -> List[str]:
    out: List[str] = []
    for ln in lines:
        m = re.match(r"^\*\*(Error\s*\d+\s*[—-]\s*[^*]+|Buena práctica[^*]*)\*\*\s*(.*)$", ln)
        if m:
            heading = _strip_md(m.group(1))
            rest = _strip_md(m.group(2))
            if rest:
                out.append(f"{heading}. {rest}")
            else:
                out.append(heading)
    return out

def _parse_actividades(lines: List[str]) -> List[Activity]:
    out: List[Activity] = []
    current: Optional[Activity] = None
    current_field: Optional[str] = None
    buffer: List[str] = []

    def flush():
        nonlocal buffer, current_field, current
        if current is None or current_field is None:
            buffer = []
            return
        v = " ".join(b.strip() for b in buffer if b.strip()).strip()
        if current_field == "descripcion":
            current.descripcion = (current.descripcion + " " + v).strip()
        elif current_field == "objetivo":
            current.objetivo = (current.objetivo + " " + v).strip()
        elif current_field == "material":
            current.material = (current.material + " " + v).strip()
        buffer = []

    for ln in lines:
        m = RE_ACTIVITY_HEADER.match(ln)
        if m:
            flush()
            if current is not None:
                out.append(current)
            current = Activity(
                title=_strip_md(m.group(1)),
                modalidad=m.group(2).strip(),
                duracion=m.group(3).strip(),
            )
            current_field = None
            continue
        if current is None:
            continue
        mf = RE_BOLD_TERM.match(ln)
        if mf:
            flush()
            label = mf.group(1).lower()
            content = mf.group(2)
            if label.startswith("descripción") or label.startswith("descripcion"):
                current_field = "descripcion"
            elif label.startswith("objetivo"):
                current_field = "objetivo"
            elif label.startswith("material"):
                current_field = "material"
            else:
                current_field = None
            if content:
                buffer.append(_strip_md(content))
            continue
        if current_field is not None and ln.strip() and not ln.strip().startswith("---"):
            buffer.append(_strip_md(ln.strip()))
    flush()
    if current is not None:
        out.append(current)
    return out

def _parse_cuestionario(lines: List[str]) -> List[QuizItem]:
    items: List[QuizItem] = []
    current: Optional[QuizItem] = None
    answer_buffer: List[str] = []
    for ln in lines:
        m = RE_QUIZ_QUESTION.match(ln)
        if m:
            if current is not None:
                if answer_buffer and not current.options:
                    current.answer = " ".join(answer_buffer).strip()
                items.append(current)
            current = QuizItem(number=int(m.group(1)), question=_strip_md(m.group(2)))
            answer_buffer = []
            continue
        if current is None:
            continue
        opt = re.match(r"^([a-d])\)\s+(.+)$", ln.strip())
        if opt:
            text = opt.group(2)
            current.options.append(f"{opt.group(1)}) {_strip_md(text)}")
            continue
        if ln.strip() and not ln.strip().startswith("---"):
            answer_buffer.append(_strip_md(ln.strip()))
    if current is not None:
        if answer_buffer and not current.options:
            current.answer = " ".join(answer_buffer).strip()
        items.append(current)
    return items

def _parse_faq(lines: List[str]) -> List[FAQItem]:
    items: List[FAQItem] = []
    current_q: Optional[str] = None
    answer_lines: List[str] = []
    for ln in lines:
        m = RE_FAQ_QUESTION.match(ln)
        if m:
            if current_q is not None:
                items.append(FAQItem(question=current_q, answer=" ".join(answer_lines).strip()))
            current_q = _strip_md(m.group(1))
            answer_lines = []
            continue
        if current_q is not None and ln.strip() and not ln.strip().startswith("---"):
            answer_lines.append(_strip_md(ln.strip()))
    if current_q is not None:
        items.append(FAQItem(question=current_q, answer=" ".join(answer_lines).strip()))
    return items

def _parse_herramientas(lines: List[str]) -> List[Tool]:
    out: List[Tool] = []
    for ln in lines:
        m = re.match(r"^-\s+\*\*([^*]+)\*\*\s*(.*)$", ln)
        if m:
            out.append(Tool(name=_strip_md(m.group(1)), description=_strip_md(m.group(2))))
    return out

def _parse_mensaje_clave(lines: List[str]) -> str:
    parts: List[str] = []
    for ln in lines:
        if ln.strip().startswith(">"):
            parts.append(_strip_md(ln.strip().lstrip(">").strip()))
    return " ".join(parts).strip().strip('"')

# ─── bank_index ──────────────────────────────────────────────────────────────

def load_bank_codes() -> set:
    if not BANK_PATH.exists():
        return set()
    text = BANK_PATH.read_text(encoding="utf-8")
    return set(re.findall(r"EJ-D\d{2}-\d+", text))

# ─── validator ───────────────────────────────────────────────────────────────

def validate(brief: Brief, bank: set, strict_hilo: bool = False) -> ValidationReport:
    report = ValidationReport(brief_code=brief.code)
    missing = SECCIONES_OBLIGATORIAS - brief.sections_found
    for n in sorted(missing):
        report.errors.append(f"Falta sección obligatoria §{n}")
    if not brief.title:
        report.errors.append("Falta título de sesión (## Módulo … / ## Sesión …)")
    if not brief.estructura_tabla:
        report.errors.append("§4 vacío o sin tabla parseable")
    if not brief.conceptos:
        report.errors.append("§6 sin ningún ### concepto")
    if not brief.ejemplos:
        report.errors.append("§8 sin ningún ### Caso")
    if not brief.actividades:
        report.errors.append("§11 sin ninguna actividad")

    if brief.code in DIAS_CON_HILO_PROYECTO and 18 not in brief.sections_found:
        msg = f"Falta §18 🎓 HILO DEL PROYECTO FINAL en {brief.code}"
        if strict_hilo:
            report.errors.append(msg)
        else:
            report.warnings.append(msg)

    if len(brief.conceptos) < 8:
        report.warnings.append(f"§6 tiene {len(brief.conceptos)} conceptos (mínimo 8)")
    for c in brief.conceptos:
        if not c.bullets_hint:
            report.warnings.append(f"Concepto '{c.title}' sin **Puntos clave:** (fallback aplicado)")
        for b in c.bullets_hint:
            if len(b.split()) > MAX_PALABRAS_POR_BULLET:
                report.warnings.append(
                    f"Bullet largo ({len(b.split())} palabras) en concepto '{c.title}'"
                )
    if len(brief.datos) < 3:
        report.warnings.append(f"§7 tiene {len(brief.datos)} cifras (mínimo 3)")
    if len(brief.ejemplos) < 4:
        report.warnings.append(f"§8 tiene {len(brief.ejemplos)} ejemplos (mínimo 4)")
    for ex in brief.ejemplos:
        if not all([ex.situacion or True, ex.prompt or ex.resultado]):
            report.warnings.append(f"Ejemplo '{ex.title}' sin todos los campos (Situación/Prompt/Resultado/Iteración)")
    if len(brief.glosario) < 6:
        report.warnings.append(f"§9 tiene {len(brief.glosario)} términos (mínimo 6)")
    if len(brief.actividades) < 3:
        report.warnings.append(f"§11 tiene {len(brief.actividades)} actividades (mínimo 3)")

    if bank:
        codes_in_brief = set(re.findall(r"EJ-D\d{2}-\d+", "\n".join(brief.ejercicios)))
        huerfanos = codes_in_brief - bank
        for c in huerfanos:
            report.warnings.append(f"Código de ejercicio huérfano: {c} no está en Banco_de_Ejercicios.md")

    return report

# ─── planner ─────────────────────────────────────────────────────────────────

def _first_sentence(text: str) -> str:
    if not text:
        return ""
    s = re.split(r"(?<=[.!?])\s+", text.strip(), maxsplit=1)
    return s[0] if s else text

def _bullets_for_concept(c: Concept) -> Tuple[List[str], str]:
    """Devuelve (bullets, notes_text). Si hay Puntos clave, los usa literal.
    Si no, fallback: primera oración de cada párrafo."""
    notes = "\n\n".join(c.paragraphs).strip()
    if c.bullets_hint:
        bullets = c.bullets_hint[:MAX_BULLETS_POR_SLIDE]
        return bullets, notes
    fallback: List[str] = []
    for p in c.paragraphs[:MAX_BULLETS_POR_SLIDE]:
        s = _first_sentence(p)
        if s:
            fallback.append(s)
    return fallback, notes

def build_slides(brief: Brief) -> List[Slide]:
    slides: List[Slide] = []

    # 1. Portada
    slides.append(Slide(
        layout="portada",
        title=brief.title or f"{brief.code} · IA en la Oficina",
        subtitle=brief.fecha,
        notes=brief.contexto,
    ))

    # 2. Agenda (§4 — tabla de bloques)
    if brief.estructura_tabla:
        slides.append(Slide(
            layout="agenda",
            title="Agenda de la sesión",
            body_text=brief.estructura_tabla,
            notes=brief.contexto,
        ))

    # 3. Recap sesión anterior (§2 — sólo si hay)
    if brief.recap_anterior and brief.code != "D01":
        slides.append(Slide(
            layout="recap",
            title="De dónde venimos",
            bullets=[brief.recap_anterior],
            notes=brief.contexto,
        ))

    # 4. Objetivos (§5)
    if brief.objetivos:
        slides.append(Slide(
            layout="objetivos",
            title="Objetivos de hoy",
            bullets=brief.objetivos[:MAX_BULLETS_POR_SLIDE],
            notes="\n".join(brief.objetivos),
        ))

    # 5. Conceptos de §6 (una slide por concepto)
    for c in brief.conceptos:
        bullets, notes = _bullets_for_concept(c)
        slides.append(Slide(
            layout="concepto",
            title=c.title,
            bullets=bullets,
            notes=notes,
        ))

    # 6. Datos y cifras (§7)
    if brief.datos:
        # Partir en varias slides si excede MAX_BULLETS
        for chunk in _chunk(brief.datos, MAX_BULLETS_POR_SLIDE):
            slides.append(Slide(
                layout="datos",
                title="Datos y cifras clave",
                bullets=chunk,
                notes="\n- " + "\n- ".join(brief.datos),
            ))

    # 7. Ejemplos resueltos (§8 — una slide por caso)
    for ex in brief.ejemplos:
        bullets: List[str] = []
        if ex.situacion:
            bullets.append(f"Situación: {ex.situacion[:120]}")
        if ex.prompt:
            bullets.append(f"Prompt: {ex.prompt[:120]}")
        if ex.resultado:
            bullets.append(f"Resultado: {ex.resultado[:120]}")
        if ex.iteracion:
            bullets.append(f"Iteración: {ex.iteracion[:120]}")
        notes = "\n\n".join(filter(None, [
            f"Situación: {ex.situacion}" if ex.situacion else "",
            f"Prompt: {ex.prompt}" if ex.prompt else "",
            f"Resultado: {ex.resultado}" if ex.resultado else "",
            f"Iteración: {ex.iteracion}" if ex.iteracion else "",
        ]))
        slides.append(Slide(
            layout="ejemplo",
            title=ex.title,
            bullets=bullets,
            notes=notes,
        ))

    # 8. Errores comunes (§10)
    if brief.errores:
        for chunk in _chunk(brief.errores, MAX_BULLETS_POR_SLIDE):
            slides.append(Slide(
                layout="errores",
                title="Errores comunes y buenas prácticas",
                bullets=chunk,
                notes="\n- " + "\n- ".join(brief.errores),
            ))

    # 9. Glosario (§9 — agrupado de 5 en 5)
    if brief.glosario:
        for chunk_terms in _chunk(brief.glosario, 5):
            bullets = [f"{t.term}: {t.definition[:80]}" for t in chunk_terms]
            notes = "\n".join(f"{t.term}: {t.definition}" for t in chunk_terms)
            slides.append(Slide(
                layout="glosario",
                title="Glosario",
                bullets=bullets,
                notes=notes,
            ))

    # 10. Slide de práctica
    slides.append(Slide(
        layout="practica",
        title="AHORA: PRÁCTICA",
        subtitle="A partir de aquí, manos al teclado",
        notes=brief.parejas,
    ))

    # 11. Fichas de ejercicio (§11 — una por actividad)
    for act in brief.actividades:
        bullets: List[str] = []
        if act.modalidad and act.duracion:
            bullets.append(f"Modalidad: {act.modalidad.strip()} · {act.duracion.strip()}")
        if act.descripcion:
            bullets.append(act.descripcion[:200])
        if act.objetivo:
            bullets.append(f"Objetivo: {act.objetivo[:120]}")
        if act.material:
            bullets.append(f"Material: {act.material[:120]}")
        notes = "\n".join(filter(None, [
            f"Modalidad: {act.modalidad} · {act.duracion}",
            f"Descripción: {act.descripcion}",
            f"Objetivo: {act.objetivo}",
            f"Material: {act.material}",
        ]))
        slides.append(Slide(
            layout="ficha_ejercicio",
            title=act.title,
            bullets=bullets,
            notes=notes,
        ))

    # 12. Parejas (§12)
    if brief.parejas:
        slides.append(Slide(
            layout="parejas",
            title="Trabajo en parejas",
            body_text=brief.parejas[:800],
            notes=brief.parejas,
        ))

    # 13. Puesta en común (mini-entregable §17)
    if brief.mini_entregable:
        slides.append(Slide(
            layout="puesta_comun",
            title="Mini-entregable del día",
            body_text=brief.mini_entregable[:600],
            notes=brief.mini_entregable,
        ))

    # 14. Mensaje clave (§19)
    if brief.mensaje_clave:
        slides.append(Slide(
            layout="mensaje_clave",
            title="Mensaje clave",
            body_text=brief.mensaje_clave,
            notes=brief.mensaje_clave,
        ))

    # 15. Ideas para recordar (§20)
    if brief.ideas:
        slides.append(Slide(
            layout="cierre_hoy",
            title="Ideas para recordar",
            bullets=brief.ideas[:MAX_BULLETS_POR_SLIDE * 2],
            notes="\n- " + "\n- ".join(brief.ideas),
        ))

    # 16. Proyecto final (§18, si aplica)
    if brief.hilo_proyecto:
        slides.append(Slide(
            layout="proyecto_final",
            title="Hilo del proyecto final",
            body_text=brief.hilo_proyecto[:800],
            notes=brief.hilo_proyecto,
        ))

    # 17. Conexión con el resto del curso (§21)
    if brief.conexion:
        slides.append(Slide(
            layout="cierre_manana",
            title="Hasta mañana",
            body_text=brief.conexion[:600],
            notes=brief.conexion,
        ))

    return slides[:MAX_SLIDES]

def _chunk(lst, size):
    for i in range(0, len(lst), size):
        yield lst[i:i + size]

# ─── renderer ────────────────────────────────────────────────────────────────

def _import_pptx():
    try:
        from pptx import Presentation  # noqa: F401
        from pptx.util import Inches, Pt, Emu  # noqa: F401
        from pptx.dml.color import RGBColor  # noqa: F401
        from pptx.enum.shapes import MSO_SHAPE  # noqa: F401
        from pptx.enum.text import PP_ALIGN  # noqa: F401
        return True
    except ImportError:
        return False

def render_deck(slides: List[Slide], out_path: Path, day_code: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank_layout = prs.slide_layouts[6]
    N = len(slides)

    for i, sl in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        is_portada = (sl.layout == "portada")

        # Fondo
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*(COLOR_CREMA if not is_portada else COLOR_OSCURO))

        if is_portada:
            _render_portada(slide, sl, day_code, prs)
        else:
            _render_content(slide, sl, day_code, i + 1, N, prs)

        # Notas del orador
        if sl.notes:
            slide.notes_slide.notes_text_frame.text = sl.notes

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

def _render_portada(slide, sl, day_code, prs):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # Logo centrado grande
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(SLIDE_W_IN / 2 - 2.0),
            Inches(1.2),
            Inches(4.0),
        )

    # Título principal
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(SLIDE_W_IN - 2.0), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = sl.title
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*COLOR_BLANCO)

    # Subtítulo (fecha)
    if sl.subtitle:
        sb = slide.shapes.add_textbox(Inches(1.0), Inches(7.8), Inches(SLIDE_W_IN - 2.0), Inches(0.8))
        sf = sb.text_frame
        sp = sf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sr = sp.add_run()
        sr.text = sl.subtitle
        sr.font.size = Pt(22)
        sr.font.color.rgb = RGBColor(*COLOR_NARANJA)

    # Badge "Día X de 17"
    _add_badge(slide, day_code)

def _render_content(slide, sl, day_code, idx, total, prs):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    # Logo esquina superior derecha
    if LOGO_PATH.exists():
        slide.shapes.add_picture(
            str(LOGO_PATH),
            Inches(SLIDE_W_IN - 1.5),
            Inches(0.2),
            Inches(1.2),
        )

    # Badge día
    _add_badge(slide, day_code)

    # Icono + título
    icon = ICONS_BY_LAYOUT.get(sl.layout, "")
    title_text = f"{icon}  {sl.title}".strip() if icon else sl.title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(SLIDE_W_IN - 2.0), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*COLOR_OSCURO)

    # Cuerpo
    if sl.bullets:
        body = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(SLIDE_W_IN - 1.4), Inches(6.5))
        bf = body.text_frame
        bf.word_wrap = True
        for j, b in enumerate(sl.bullets):
            if j == 0:
                bp = bf.paragraphs[0]
            else:
                bp = bf.add_paragraph()
            br = bp.add_run()
            br.text = f"• {b}"
            br.font.size = Pt(22)
            br.font.color.rgb = RGBColor(*COLOR_OSCURO)
            bp.space_after = Pt(8)
    elif sl.body_text:
        body = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(SLIDE_W_IN - 1.4), Inches(6.5))
        bf = body.text_frame
        bf.word_wrap = True
        for j, line in enumerate(sl.body_text.split("\n")):
            line = line.rstrip()
            if not line:
                continue
            if j == 0:
                bp = bf.paragraphs[0]
            else:
                bp = bf.add_paragraph()
            br = bp.add_run()
            br.text = line
            br.font.size = Pt(18)
            br.font.color.rgb = RGBColor(*COLOR_OSCURO)

    # Footer + numeración
    fb = slide.shapes.add_textbox(Inches(0.5), Inches(SLIDE_H_IN - 0.45), Inches(SLIDE_W_IN - 1.5), Inches(0.35))
    ff = fb.text_frame
    fp = ff.paragraphs[0]
    fr = fp.add_run()
    fr.text = FOOTER_TEXT
    fr.font.size = Pt(10)
    fr.font.color.rgb = RGBColor(*COLOR_GRIS)

    nb = slide.shapes.add_textbox(Inches(SLIDE_W_IN - 1.5), Inches(SLIDE_H_IN - 0.45), Inches(1.0), Inches(0.35))
    nf = nb.text_frame
    np_ = nf.paragraphs[0]
    np_.alignment = PP_ALIGN.RIGHT
    nr = np_.add_run()
    nr.text = f"{idx}/{total}"
    nr.font.size = Pt(10)
    nr.font.color.rgb = RGBColor(*COLOR_GRIS)

def _add_badge(slide, day_code):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN

    day_num = int(re.sub(r"\D", "", day_code) or 0)
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5),
        Inches(0.25),
        Inches(2.0),
        Inches(0.45),
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*COLOR_AZUL)
    shp.line.color.rgb = RGBColor(*COLOR_AZUL)
    tf = shp.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"Día {day_num} de 17"
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*COLOR_BLANCO)

# ─── cli ─────────────────────────────────────────────────────────────────────

def find_brief_path(code: str) -> Optional[Path]:
    code = code.upper()
    candidates = sorted(BRIEFS_DIR.glob(f"{code}_*.md"))
    return candidates[0] if candidates else None

def _ir_to_dict(brief: Brief, slides: List[Slide]) -> dict:
    return {
        "brief": {
            "code": brief.code,
            "title": brief.title,
            "fecha": brief.fecha,
            "n_conceptos": len(brief.conceptos),
            "n_ejemplos": len(brief.ejemplos),
            "n_actividades": len(brief.actividades),
            "n_glosario": len(brief.glosario),
            "n_cuestionario": len(brief.cuestionario),
            "n_faq": len(brief.faq),
            "n_ideas": len(brief.ideas),
            "sections_found": sorted(brief.sections_found),
        },
        "slides": [
            {
                "layout": s.layout,
                "title": s.title,
                "n_bullets": len(s.bullets),
                "has_notes": bool(s.notes),
            }
            for s in slides
        ],
    }

def process_one(code: str, args, bank: set) -> int:
    path = find_brief_path(code)
    if path is None:
        print(f"❌ No se encontró brief para {code} en {BRIEFS_DIR}")
        return 3
    print(f"\n📄 {code} → {path.name}")
    brief = parse_brief(path)
    report = validate(brief, bank, strict_hilo=args.strict_hilo)
    report.print_summary()
    if report.has_errors():
        return 2

    slides = build_slides(brief)
    print(f"  📊 Slides previstas: {len(slides)}")

    if args.dump_ir:
        ir = _ir_to_dict(brief, slides)
        ir_path = Path(args.dump_ir)
        ir_path.write_text(json.dumps(ir, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"  💾 IR dumped: {ir_path}")

    if args.validate:
        return 1 if report.warnings else 0

    if not _import_pptx():
        print("❌ python-pptx no instalado. Ejecuta: pip install -r requirements.txt")
        return 3

    out_path = Path(args.out) if args.out else DECKS_DIR / f"{brief.code}_procesado.pptx"
    render_deck(slides, out_path, brief.code)
    print(f"  💾 PPTX: {out_path}")
    return 1 if report.warnings else 0

def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Generar deck del formador desde brief.")
    parser.add_argument("codes", nargs="*", help="Códigos DXX (D01, D02, ...)")
    parser.add_argument("--all", action="store_true", help="Procesar todos los DXX presentes")
    parser.add_argument("--validate", action="store_true", help="Solo validar, no generar PPTX")
    parser.add_argument("--dump-ir", metavar="PATH", help="Guardar IR como JSON para snapshot")
    parser.add_argument("--strict-hilo", action="store_true",
                        help="§18 ausente en D01/D09/D16/D17 = error fatal")
    parser.add_argument("--out", metavar="PATH", help="Override path de salida del PPTX")
    args = parser.parse_args(argv)

    codes: List[str] = []
    if args.all:
        codes = sorted({p.stem.split("_")[0] for p in BRIEFS_DIR.glob("D*_*.md")})
    else:
        codes = [c.upper() for c in args.codes]

    if not codes:
        parser.print_help()
        return 3

    bank = load_bank_codes()
    overall = 0
    for code in codes:
        rc = process_one(code, args, bank)
        overall = max(overall, rc)

    return overall

if __name__ == "__main__":
    sys.exit(main())
