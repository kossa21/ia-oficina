#!/usr/bin/env python3
"""
Pipeline completo — IA en la Oficina 01/26
Brief .md  →  NotebookLM (crea notebook + genera deck)  →  PPTX con logo y notas

Uso:
  python pipeline_ia_oficina.py           # procesa todas las sesiones pendientes
  python pipeline_ia_oficina.py D05 D06   # procesa solo las sesiones indicadas
"""

import asyncio
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# ──────────────────────────────────────────────────────────
# RUTAS
# ──────────────────────────────────────────────────────────

BASE = Path("/Users/ananda/Desktop/Ananda/INSERTA - IA en la oficina (shared space)/Edicion 01-26")
LOGO = BASE / "Logos/00_FULL COLOR (1).png"

# ──────────────────────────────────────────────────────────
# PROMPT ESTÁNDAR para NotebookLM
# ──────────────────────────────────────────────────────────

PROMPT = (
    "Genera una presentación visualmente rica basada en el contenido de este brief. "
    "Usa títulos provocadores y narrativos, no descriptivos. "
    "Cada slide debe tener un elemento visual (diagrama, estadística, comparativa o iconos). "
    "Estructura obligatoria del deck, en este orden exacto: "
    "(1) PORTADA con el título de la sesión, 'Día X de 17', módulo y nombre del curso. "
    "(2) AGENDA DEL DÍA: una slide con los bloques de la sesión y su duración "
    "(Apertura 15 min, Teórico-demostrativo 75 min, Descanso 15 min, "
    "Práctica guiada 95 min, Puesta en común 40 min) y, junto a cada bloque, "
    "los temas o actividades que se verán. "
    "(3) SLIDES DE CONTENIDO: máx. 5 puntos cortos por slide. "
    "(4) SEÑAL DE EJERCICIO: justo antes de cada actividad práctica, una slide de "
    "transición '🛠 AHORA: PRÁCTICA' con el nombre de la actividad, su código "
    "(EJ-DXX-N si aparece en el documento), la modalidad (Individual/En parejas/En grupo), "
    "la duración en minutos y el objetivo en una frase, para que la formadora sepa "
    "cuándo parar la teoría y hacer la práctica. "
    "(5) CIERRE: una slide 'Lo que hicimos hoy' con 3-4 puntos de recap y una slide "
    "'Mañana' con 2-3 puntos de la sesión siguiente. "
    "Tono: directo, accesible, sin tecnicismos. "
    "Público: personal administrativo sin experiencia previa en IA."
)

# ──────────────────────────────────────────────────────────
# SESIONES — añade aquí D10–D17 cuando tengas los briefs
# ──────────────────────────────────────────────────────────

SESIONES = [
    # (brief_md,                      notebook_name,                                  codigo)
    ("D05_Brief_NotebookLM.md",  "IA Oficina 01/26 — D05 Resúmenes y Formatos",  "D05"),
    ("D06_Brief_NotebookLM.md",  "IA Oficina 01/26 — D06 Correos Profesionales", "D06"),
    ("D07_Brief_NotebookLM.md",  "IA Oficina 01/26 — D07 Informes con IA",       "D07"),
    ("D08_Brief_NotebookLM.md",  "IA Oficina 01/26 — D08 Planificación Semanal", "D08"),
    ("D09_Brief_NotebookLM.md",  "IA Oficina 01/26 — D09 Actas y Delegación",    "D09"),
    ("D10_Brief_NotebookLM.md", "IA Oficina 01/26 — D10 Análisis de Datos",     "D10"),
    ("D11_Brief_NotebookLM.md", "IA Oficina 01/26 — D11 Búsqueda e Informes",   "D11"),
    ("D12_Brief_NotebookLM.md", "IA Oficina 01/26 — D12 Canva con IA",          "D12"),
    ("D13_Brief_NotebookLM.md", "IA Oficina 01/26 — D13 Comunicación Visual",   "D13"),
    ("D14_Brief_NotebookLM.md", "IA Oficina 01/26 — D14 Gestión Documental",    "D14"),
    ("D15_Brief_NotebookLM.md", "IA Oficina 01/26 — D15 Atención al Cliente",   "D15"),
    ("D16_Brief_NotebookLM.md", "IA Oficina 01/26 — D16 Proyecto Personal",     "D16"),
    ("D17_Brief_NotebookLM.md", "IA Oficina 01/26 — D17 Presentaciones Finales","D17"),
]

# ──────────────────────────────────────────────────────────
# FASE 1 — Generar deck en NotebookLM y descargar PPTX
# ──────────────────────────────────────────────────────────

async def generar_deck(client, brief_md, notebook_name, codigo):
    raw_path = BASE / "Decks" / f"{codigo}_raw.pptx"

    if raw_path.exists():
        print(f"  ⏭  {codigo}: PPTX raw ya existe, saltando generación")
        return raw_path

    print(f"\n  📓 Creando notebook: {notebook_name}")
    nb = await client.notebooks.create(notebook_name)

    print(f"  📎 Subiendo brief: {brief_md}")
    await client.sources.add_file(nb.id, str(BASE / "Briefs" / brief_md), wait=True)

    print(f"  🎨 Generando presentación (1–3 min)...")
    status = await client.artifacts.generate_slide_deck(nb.id, instructions=PROMPT)
    await client.artifacts.wait_for_completion(nb.id, status.task_id)

    print(f"  ⬇  Descargando PPTX...")
    await client.artifacts.download_slide_deck(nb.id, str(raw_path), output_format="pptx")

    print(f"  ✅ Guardado: {raw_path.name}")
    return raw_path


# ──────────────────────────────────────────────────────────
# FASE 2 — Añadir logo AllWomen y speaker notes al PPTX
# ──────────────────────────────────────────────────────────

def leer_brief(brief_md):
    """Extrae líneas clave del brief para usar como base de las notas."""
    path = BASE / "Briefs" / brief_md
    if not path.exists():
        return []
    lines = path.read_text(encoding="utf-8").splitlines()
    # Devuelve los bloques de contenido (líneas no vacías y no de cabecera)
    return [l.strip() for l in lines if l.strip() and not l.startswith("#") and not l.startswith("|") and not l.startswith("-")]


def nota_para_slide(i, lineas_brief, total_slides):
    """Genera una nota contextual para cada slide."""
    if i == 0:
        return "PORTADA — Presenta el tema del día y conecta con lo que el grupo aprendió en la sesión anterior. Recuerda mencionar el mini-entregable al final."
    if i == total_slides - 1:
        return "CIERRE — Resume el mensaje clave de la sesión. Presenta el mini-entregable y anticipa brevemente el contenido de la próxima sesión."
    # Para slides intermedios, rota por las líneas del brief
    idx = ((i - 1) * 3) % max(len(lineas_brief), 1)
    fragmento = " | ".join(lineas_brief[idx:idx+3]) if lineas_brief else "Desarrolla el contenido con ejemplos del puesto de trabajo del alumno."
    return f"Slide {i+1}: {fragmento[:300]}"


def procesar_pptx(raw_path, codigo, brief_md):
    output_path = BASE / "Decks" / f"{codigo}_procesado.pptx"

    prs = Presentation(str(raw_path))

    # Dimensiones del logo
    img = Image.open(LOGO)
    logo_w_in = 1.2
    logo_h_in = logo_w_in / (img.size[0] / img.size[1])

    slide_w_in = prs.slide_width.inches
    left = Inches(slide_w_in - logo_w_in - 0.3)
    top = Inches(0.2)

    lineas_brief = leer_brief(brief_md)
    total = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        # Logo
        slide.shapes.add_picture(
            str(LOGO), left, top,
            width=Inches(logo_w_in),
            height=Inches(logo_h_in)
        )
        # Nota
        nota = nota_para_slide(i, lineas_brief, total)
        slide.notes_slide.notes_text_frame.text = nota

    prs.save(str(output_path))
    print(f"  ✅ {output_path.name}  ({total} slides, logo + notas añadidos)")
    return output_path


# ──────────────────────────────────────────────────────────
# MAIN
# ──────────────────────────────────────────────────────────

async def main():
    from notebooklm import NotebookLMClient

    print()
    print("=" * 56)
    print("  Pipeline IA en la Oficina 01/26")
    print("  Brief → NotebookLM → PPTX con logo y notas")
    print("=" * 56)

    # Filtrar sesiones si se pasan argumentos
    sesiones = SESIONES
    if len(sys.argv) > 1:
        codigos = [a.upper() for a in sys.argv[1:]]
        sesiones = [s for s in SESIONES if s[2] in codigos]
        if not sesiones:
            print(f"\n⚠️  No se encontraron sesiones para: {sys.argv[1:]}")
            print(f"   Sesiones disponibles: {[s[2] for s in SESIONES]}")
            return
        print(f"\nProcesando: {[s[2] for s in sesiones]}")

    raw_resultados = []

    # ── Fase 1: NotebookLM ──
    print("\n── Fase 1: Generación en NotebookLM ─────────────\n")
    async with await NotebookLMClient.from_storage() as client:
        for brief_md, notebook_name, codigo in sesiones:
            try:
                raw = await generar_deck(client, brief_md, notebook_name, codigo)
                raw_resultados.append((raw, codigo, brief_md))
            except Exception as e:
                print(f"  ❌ Error en {codigo}: {e}")

    # ── Fase 2: Logo + notas ──
    print("\n── Fase 2: Logo AllWomen + speaker notes ─────────\n")
    procesados = []
    for raw_path, codigo, brief_md in raw_resultados:
        try:
            out = procesar_pptx(raw_path, codigo, brief_md)
            procesados.append(out)
        except Exception as e:
            print(f"  ❌ Error procesando {codigo}: {e}")

    # ── Resumen ──
    print()
    print("=" * 56)
    print(f"  ✅ Completados: {len(procesados)}/{len(sesiones)}")
    for p in procesados:
        print(f"     {p.name}")
    print("=" * 56)
    print()


asyncio.run(main())
