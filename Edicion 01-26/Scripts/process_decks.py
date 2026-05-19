#!/usr/bin/env python3
"""
process_decks.py — IA en la Oficina 01/26
Añade logo AllWomen + speaker notes a los PPTX generados por NotebookLM.

Uso:
    python3 process_decks.py D05_MiDeck.pptx
    python3 process_decks.py D05_MiDeck.pptx D06_Otro.pptx D07_Otro.pptx
    python3 process_decks.py *.pptx   (todos los pptx sin procesar)

El script detecta automáticamente el día (D05, D06...) desde el nombre del archivo
y aplica las speaker notes correspondientes.
El archivo procesado se guarda en el mismo directorio con el nombre original.
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ── CONFIGURACIÓN ────────────────────────────────────────────────────────────

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(SCRIPT_DIR, "..", "Logos", "00_FULL COLOR (1).png")

# Posición logo AllWomen (aprobada por Daniela, no cambiar)
LOGO_WIDTH  = Inches(1.2)
LOGO_TOP    = Inches(0.20)
SLIDE_W     = Inches(17.78)
LOGO_LEFT   = SLIDE_W - LOGO_WIDTH - Inches(0.30)

# ── SPEAKER NOTES POR SESIÓN ─────────────────────────────────────────────────
# Cada lista tiene una nota por slide, en orden.
# Si el deck tiene más slides que notas, las slides extra quedan sin nota.
# Si tiene menos, solo se usan las primeras.

NOTES = {

    "D05": [
        "Arranca el M3 sesión 2. Conecta con D04: ya automatizan tareas repetitivas. Hoy aprenden a transformar cualquier documento en el formato exacto que necesita cada destinatario.",
        "Lanza la pregunta: '¿Alguna vez habéis tenido que explicar lo mismo tres veces a tres personas distintas?' Ese problema lo resuelve la IA en 30 segundos.",
        "Los tres tipos de resumen son el núcleo de la sesión. Escribe en la pizarra: EJECUTIVO / OPERATIVO / DIVULGATIVO. Pregunta: ¿Para quién escribiría cada uno de vosotros normalmente?",
        "Demuestra en directo: coge el acta ficticia del material, pégala en ChatGPT, y genera los tres resúmenes uno a uno con el prompt exacto de la slide. El grupo ve la diferencia en tiempo real.",
        "Conversión de formatos: el más útil en el día a día es texto→tabla y tabla→texto. Demuestra ambos con el mismo documento que acabas de usar. Pregunta: '¿Cuál usaríais vosotros esta semana?'",
        "Extracción de campos: este es el que más impresiona. Pega un correo largo y extrae: fechas, importes, responsables, acciones. El grupo ve que lo que antes tomaba 10 minutos ahora tarda 20 segundos.",
        "Documentos largos: no entres en tecnicismos sobre límites de tokens. Explica simplemente: 'Si el documento es muy largo, lo partes en trozos y luego pides el resumen global.'",
        "Actividad 1 — Acta en tres formatos (40 min). Proporciona el acta ficticia del material. El grupo trabaja individualmente y comparte resultados. Señala las diferencias de vocabulario entre versiones.",
        "Actividad 2 — Tu documento real (50 min). Cada alumno procesa un documento propio anonimizado. Circula por las mesas, ayuda con los prompts. Es el momento de personalizar el aprendizaje.",
        "Cierre y mini-entregable: tres versiones del mismo documento. Pregunta final: '¿Qué resumen usarías mañana en el trabajo?' Que cada persona nombre su destinatario real.",
    ],

    "D06": [
        "Arranca M3 sesión 3. Conecta con D05: ya saben resumir y transformar documentos. Hoy todo gira alrededor del correo, que es donde la mayoría pasa entre 1 y 2 horas al día.",
        "Estadística de impacto: el 47% de los correos no reciben respuesta en 24h por un asunto o cierre mal redactado. Pregunta: '¿Cuántos de vosotros tenéis correos sin respuesta desde hace más de 48h?'",
        "Los 6 elementos del correo: repásalos uno a uno con ejemplos malos vs buenos. El asunto es clave: demuestra el mismo correo con asunto vago ('Reunión') vs asunto concreto ('Validación del informe Q2 — plazo viernes').",
        "Situaciones difíciles: estas son las que más valor tienen para el grupo. Haz una votación rápida: '¿Cuál de estas situaciones os genera más bloqueo?' Y empieza por la más votada.",
        "Demuestra en directo: coge la situación más votada, escribe el prompt de la slide, y genera el correo en tiempo real. Pide al grupo que sugiera mejoras antes de cerrar.",
        "Respuestas a clientes — el patrón de 3 pasos: RECONOCE / INFORMA / ACTÚA. Escríbelo en la pizarra. Es la estructura que deben memorizar. Úsala en el role-play.",
        "Revisión con IA antes de enviar: demuestra el prompt de auditoría con un correo real (tuyo, anonimizado). El grupo ve cómo la IA señala lo que no está claro.",
        "Actividad 1 — 3 correos difíciles en directo (25 min). El grupo propone las situaciones, tú las resuelves con IA en pantalla. Mantén el ritmo: máximo 8 minutos por correo.",
        "Actividad 2 — Mis 5 correos reales (55 min). Cada alumno reescribe 5 correos propios. Insiste: deben anonimizar nombres y datos. Circula por las mesas y señala mejoras concretas.",
        "Role-play y mini-entregable: el banco de 5 plantillas. Que cada alumno comparta una plantilla antes de cerrar. Recuerda: los [CORCHETES] son los campos variables que deberán personalizar cada vez.",
    ],

    "D07": [
        "Arranca M3 sesión 4 y cierre del módulo de redacción. Conecta con D06: ya tienen su banco de plantillas. Hoy construyen documentos más largos: el informe.",
        "El informe es el mayor bloqueador de productividad en oficina. Pregunta: '¿Cuánto tiempo os lleva redactar un informe de 2 páginas?' Anota las respuestas. Al final de la sesión lo compararán con el tiempo que tardaron hoy.",
        "El flujo de 4 pasos es el contenido más importante de la sesión. Escríbelo en la pizarra: ÍNDICE → SECCIONES → RESUMEN EJECUTIVO → REVISIÓN. Insiste: el resumen ejecutivo siempre al final, nunca al principio.",
        "Demuestra paso a paso: usa el caso ficticio del material, construye el informe en 4 prompts en tiempo real. El grupo ve que el resultado en 15 minutos equivale a lo que normalmente les lleva una tarde.",
        "Adaptación de tono — las 3 versiones: TÉCNICA / EJECUTIVA / CLIENTE. Muestra las diferencias en el mismo párrafo. El vocabulario, la longitud de las frases y los detalles cambian radicalmente.",
        "Alucinaciones en informes: este es el punto de seguridad más crítico de la sesión. La regla de oro: 'Si no lo escribiste tú, verifícalo.' Demuestra el prompt de auditoría: pide a la IA que liste todos los datos que usó.",
        "Actividad 1 — Informe en 4 iteraciones (55 min). Proporciona los datos brutos del caso. El grupo construye el informe colectivamente, tú guías cada prompt. Al final, cronometrad el tiempo total.",
        "Actividad 2 — Adaptar a 3 públicos (40 min). Cada alumno adapta un texto propio. Sugiere que elijan un informe o documento que tengan que entregar próximamente — el entregable será real.",
        "Revisión cruzada y mini-entregable (10 min). Las 3 versiones del informe. Pregunta final: '¿Cuánto habríais tardado en hacer esto sin IA?' Cierra el módulo M3 recordando lo que han aprendido en D04–D07.",
    ],

    "D08": [
        "Arranca el módulo M4 — Organización del trabajo. Conecta: el grupo ya domina la redacción. Ahora IA para pensar, priorizar y organizarse, no solo para escribir.",
        "Dato de impacto: los profesionales que priorizan formalmente completan un 25% más de trabajo de alto impacto. Pregunta: '¿Cómo priorizáis ahora mismo? ¿Por orden de llegada, por urgencia, por intuición?'",
        "Eisenhower con IA: demuestra en directo. Pide al grupo que dicten 8-10 tareas reales. Pégalas en ChatGPT con el prompt de la slide. El grupo ve su propia semana clasificada en tiempo real.",
        "MoSCoW: explica que es más útil para proyectos con plazo fijo. Ejemplo: preparar el curso en sí mismo. ¿Qué es Must? ¿Qué es Won't? El grupo lo reconoce porque están viviendo la planificación del curso.",
        "Planificación semanal — el método de 4 pasos: VUELCA / RESTRICCIONES / PLAN / REFINA. El 'refina' es clave: la IA no es rígida, puedes negociar con ella igual que con un asistente.",
        "Demuestra el prompt de planificación semanal con la semana ficticia del material (15 tareas). El grupo ve cómo la IA propone bloques horarios y agrupa tareas similares.",
        "Reuniones: orden del día, briefing previo, preguntas clave. Muestra cómo preparar una reunión en 5 minutos con IA vs los 20-30 minutos habituales. Pregunta: '¿Cuántas reuniones tenéis esta semana sin orden del día?'",
        "Transcripciones y privacidad: punto crítico. Antes de grabar cualquier reunión, consentimiento. Antes de subir la transcripción a ChatGPT, verificar que no hay datos confidenciales.",
        "Actividad 1 — Semana ficticia (30 min). Proporciona las 15 tareas con estimaciones. Trabajan colectivamente. Circula y ayuda con las restricciones ('tengo reunión fija el martes a las 16h').",
        "Actividad 2 — Mi próxima semana real (50 min). Este es el entregable que se llevan. Insiste en que usen tareas reales. Al final, pide voluntarios que compartan su plan.",
        "Notas → acta (25 min) y mini-entregable: el plan semanal. Que cada alumno comparta una decisión de priorización que les sorprendió hacer con IA.",
    ],

    "D09": [
        "Arranca M4 sesión 2 y cierre del módulo de organización. Conecta con D08: ya planifican y preparan reuniones. Hoy cierran el ciclo: qué pasa después de la reunión.",
        "Dato de impacto: el 60% de los acuerdos tomados en reuniones sin acta no se ejecutan en el plazo acordado. Con acta estructurada, baja al 25%. Pregunta: '¿Cuántas de vuestras reuniones tienen acta?' Normalmente la respuesta es pocas.",
        "Anatomía del acta útil: no es una transcripción, es un documento de acción. Escribe en la pizarra: DECISIONES + ACCIONES (tabla) + PRÓXIMA REUNIÓN. Eso es todo lo que importa.",
        "Demuestra el prompt de acta: pega las notas crudas del material y genera el acta en tiempo real. El grupo ve la diferencia entre el caos de las notas y la claridad del resultado.",
        "Tablero kanban en texto: no necesitan Trello ni Asana. ChatGPT genera una tabla kanban con columnas Pendiente / En curso / Bloqueado / Completado. Demuestra con el proyecto ficticio del material.",
        "Sistema RAG (🔴🟡🟢): explica que es el estándar internacional para semáforos de seguimiento. Muestra cómo la IA rellena el estado inicial y cómo el equipo lo actualiza manualmente.",
        "Instrucciones de delegación — los 5 elementos: QUÉ / POR QUÉ / CUÁNDO / CÓMO / CÓMO INFORMAR. El 'por qué' es el más olvidado y el más importante: sin contexto, la persona delegada no puede tomar buenas decisiones ante imprevistos.",
        "Recordatorios: muestra los dos prompts (primer seguimiento / segundo seguimiento sin respuesta). La diferencia de tono entre ambos es sutil pero importante. Pregunta: '¿Tenéis tareas delegadas que llevan semanas sin respuesta?'",
        "Actividad 1 — Transcripción → acta + plan de acción (40 min). Proporciona la transcripción ficticia. El grupo trabaja en grupos de 2-3. Compara las actas resultantes: ¿coinciden en las decisiones y responsables?",
        "Actividad 2 — Mi kanban (35 min). Cada alumno crea un tablero para una iniciativa real. Circula y ayuda a definir qué va en cada columna.",
        "Actividad 3 — 3 delegaciones reales (30 min) y mini-entregable: el kit de organización (acta + kanban + delegación). Cierra M4: '¿Qué vais a hacer diferente la próxima semana?'",
    ],

    "D10": [
        "Arranca el módulo M5 — Análisis de datos. Conecta con M4: ya organizan y delegan. Ahora aprenden a extraer valor de sus propias hojas de datos sin necesidad de ser expertos en fórmulas.",
        "Dato de impacto: el 80% del tiempo dedicado a análisis de datos en oficinas se va en limpiar y formatear, no en interpretar. La IA invierte esa proporción. Pregunta: '¿Cuánto tardáis en sacar conclusiones de una hoja de datos?'",
        "Principio de anonimización antes de análisis: SIEMPRE. Muestra en tiempo real cómo cambiar nombres por CLIENTE_A, CLIENTE_B y cifras por DATOS_VENTAS. El grupo ve que el análisis funciona igual con datos ficticios.",
        "Google Sheets + Gemini: demuestra cómo escribir en lenguaje natural lo que necesitas. 'Añade una columna que calcule el porcentaje de variación entre enero y febrero' → Gemini escribe la fórmula. El grupo ve que no hace falta saber la sintaxis.",
        "Hallazgos en lenguaje natural: copia la tabla (anonimizada) en ChatGPT y pide 3 hallazgos. Muestra cómo la IA identifica patrones, anomalías y tendencias. Insiste en la regla de verificación: comprueba manualmente el dato más llamativo.",
        "Gráficos: explica qué tipo de gráfico va con qué tipo de dato (barras=comparativas, línea=evolución, dispersión=correlaciones). Demuestra con el prompt 'Sugiere el gráfico más adecuado para estos datos y explica por qué'.",
        "Alucinaciones numéricas: este es el punto de seguridad más crítico de la sesión. La IA puede inventar números. Muestra el prompt de verificación: 'Lista todos los datos concretos que has usado para llegar a estas conclusiones'. El alumno coteja con la fuente.",
        "Actividad 1 — Demo: tabla → 3 hallazgos + gráfico (demo formador, 15 min). El grupo ve el flujo completo de principio a fin antes de practicarlo.",
        "Actividad 2 — Dataset ficticio en parejas (50 min): 5 preguntas de negocio sobre un dataset de ejemplo. Circula, ayuda con los prompts de anonimización y con la interpretación de los hallazgos.",
        "Actividad 3 — Hoja propia anonimizada (30 min individual). El entregable real: mini-análisis con 5 hallazgos + 1 gráfico + 1 recomendación. Que cada alumno comparta el hallazgo más útil que encontró.",
    ],

    "D11": [
        "Arranca M5 sesión 2. Conecta con D10: ya analizan datos propios. Hoy aprenden a buscar e investigar con IA para sintetizar información externa y producir informes de mercado o proveedores.",
        "El problema de la búsqueda sin IA: miles de resultados, tiempo perdido triangulando fuentes, difícil comparar opciones. La IA actúa como un investigador junior: hace el barrido inicial y presenta las opciones organizadas.",
        "Perplexity vs Google: muestra en tiempo real la misma búsqueda en ambas herramientas. Perplexity cita sus fuentes; Google devuelve links. Para investigación profesional, Perplexity ahorra pasos. Ambas son gratuitas.",
        "Verificación de fuentes: regla del día. Antes de incluir un dato en un informe, verifica la fuente original. Demuestra el prompt 'Cita las fuentes exactas de este dato' y muestra cómo la IA a veces admite que no puede verificarlo.",
        "Síntesis comparativa: el prompt de comparativa de proveedores con tabla de criterios (precio, calidad, plazo, referencias). Demuestra que en 3 prompts se genera lo que antes llevaba 2 horas de búsqueda manual.",
        "NotebookLM para investigación documental: sube 2-3 PDFs o artículos y haz preguntas al notebook. Muestra cómo la IA responde basándose solo en los documentos cargados, sin inventar datos externos.",
        "Estructura del informe de mercado: portada + resumen ejecutivo + tabla comparativa + fuentes + recomendación. Demuestra con el prompt de 'informe de 2 páginas sobre [tema del grupo]'.",
        "Actividad 1 — Demo investigación de proveedores (formador, 15 min): búsqueda → síntesis → tabla → conclusión. El grupo ve el flujo completo.",
        "Actividad 2 — Grupo de 3: comparativa de 3 opciones reales (50 min). Cada grupo elige un tema de su sector y produce la tabla de criterios con fuentes verificadas. El formador circula y ayuda con los prompts de verificación.",
        "Actividad 3 — Mini-informe individual (30 min) y mini-entregable: informe de mercado de 1-2 págs con fuentes verificadas. Que cada alumno comparta la fuente más útil que encontró.",
    ],

    "D12": [
        "Arranca el módulo M6 — Creatividad y diseño. Conecta con M5: ya analizan e investigan. Ahora aprenden a comunicar visualmente lo que saben. Canva es la herramienta: gratis, sin instalación, con IA integrada.",
        "El miedo al diseño: '¿No es eso para diseñadores?' Ese es el mito a romper. Los 4 principios básicos (contraste, alineación, repetición, proximidad) funcionan incluso sin saber nada de diseño. Son las reglas que Canva aplica automáticamente con Magic Design.",
        "Magic Design: demuestra el flujo completo. Selecciona un tema, sube una imagen o escribe el contenido, y Magic Design propone 8 plantillas en 10 segundos. El grupo elige la mejor y la personaliza.",
        "Magic Write: el asistente de texto dentro de Canva. Demuestra: selecciona un cuadro de texto vacío → Magic Write → escribe el brief del contenido. Canva redacta el texto en el estilo de la plantilla.",
        "Generación de imágenes con IA: cuotas del plan gratuito (varían; actualmente ~50 generaciones/mes). Muestra el prompt de imagen con estilo profesional. Anticipa: si la cuota se acaba, usar imágenes de la biblioteca de Canva (gratis y sin derechos).",
        "Derechos de uso: las imágenes generadas con IA en Canva free son de uso personal y comercial, pero verificar siempre las condiciones actuales. Para uso en empresa: avisar y revisar los términos de Canva.",
        "Actividad 1 — Demo Canva Magic completo (formador, 15 min): de brief de texto a 5 slides con imagen en tiempo real. El grupo ve que no hace falta talento artístico.",
        "Actividad 2 — En parejas: rehacer slide fea → buena (50 min). Proporciona una slide de ejemplo con mala maquetación. La pareja la rediseña aplicando los 4 principios. Comparan antes/después.",
        "Actividad 3 — Individual: 5 slides de un tema de su puesto (30 min) y mini-entregable: 5 slides + 1 imagen IA. Que cada alumno muestre su slide favorita y explique qué principio aplicó.",
    ],

    "D13": [
        "Arranca M6 sesión 2. Conecta con D12: ya saben diseñar slides. Hoy crean piezas de comunicación interna: carteles, infografías, posts de Teams, banners. Y aprenden a hacerlo accesible — crítico para Fundación ONCE.",
        "Accesibilidad: este es el punto diferencial del día. Explicar que el diseño accesible no es un añadido, es mejor diseño para todos. Las 4 reglas: contraste mínimo 4.5:1, texto alternativo en imágenes, fuente mínima 16px, nunca el color como único diferenciador.",
        "Alt text en Canva: demuestra cómo añadir texto alternativo a cada imagen. Paso a paso: selecciona imagen → clic derecho → 'Alt text' → escribe descripción concisa. El grupo lo hace con sus diseños de D12.",
        "Contraste: usa el verificador online de contraste (WebAIM Contrast Checker, gratuito). Demuestra: copia el hex del color del texto y del fondo, pega en el verificador. El grupo ve si pasa o no el ratio 4.5:1.",
        "Kit de comunicación interna: no son piezas sueltas, son una mini-campaña coherente. Misma paleta, mismo logotipo, mismo tono. Demuestra cómo crear un kit de marca en Canva: colores guardados, fuentes guardadas, plantilla reutilizable.",
        "Video simple en Canva: usa una presentación de Canva y activa la opción de grabación de presentación (gratuita). 1-2 minutos de narración sobre las slides = vídeo de comunicación interna sin editar.",
        "Actividad 1 — Demo infografía desde datos de D10 (formador, 15 min): los datos analizados el martes se convierten en una infografía visual. El grupo ve la conexión entre módulos.",
        "Actividad 2 — Grupo: mini-campaña interna de 3 piezas coherentes (50 min): cartel + post de Teams + banner de email. El grupo elige el tema (bienvenida a nuevos compañeros, aviso de proceso, celebración, etc.).",
        "Actividad 3 — Individual: 1 pieza accesible de su puesto (30 min) y mini-entregable: kit de 2-3 piezas visuales. Que cada alumno compruebe el contraste de su pieza con el verificador antes de entregar.",
    ],

    "D14": [
        "Arranca el módulo M7 — Gestión documental. Conecta con M6: ya comunican visualmente. Ahora ordenan los procesos por escrito: naming, SOPs, checklists, flujos. El caos documental tiene solución.",
        "El coste del caos documental: se estima que un trabajador de oficina pierde entre 30 min y 1h al día buscando documentos o repitiendo trabajo ya hecho. Pregunta: '¿Tenéis archivos que empiezan por 'final', 'final2', 'definitivo'?' El grupo sonríe.",
        "Naming convention con IA: demuestra el prompt de nomenclatura de archivos. Proporciona una carpeta ficticia caótica y el resultado estandarizado. El grupo ve la diferencia en 2 minutos.",
        "SOP en 3 prompts: contexto del proceso → pasos detallados → validación ('¿Qué podría salir mal?'). Demuestra con el proceso de alta de un nuevo proveedor. El resultado es un documento que cualquiera puede seguir.",
        "Checklist desde el SOP: un prompt más convierte el SOP en una lista de verificación accionable. Muestra cómo la IA transforma las instrucciones en casillas de verificación ordenadas por fase.",
        "Mapeo de flujo en texto: explica que no hace falta un diagrama de Visio para mapear un proceso. Un texto con INICIO → paso 1 → decisión (¿sí/no?) → paso 2 → FIN es suficiente para identificar cuellos de botella. La IA genera ese mapa desde una descripción libre.",
        "Google Forms → Sheets → IA: demuestra el flujo: formulario de recepción de solicitudes → respuestas en Sheets → análisis con IA. Sin código, todo gratis con Gmail.",
        "NotebookLM como base de conocimiento del equipo: sube los SOPs del equipo al notebook, comparte el enlace, y cualquier persona puede preguntar '¿Cuál es el proceso para X?' La IA responde basándose solo en los documentos del equipo.",
        "Actividad 1 — Demo proceso caótico → SOP + checklist (formador, 15 min). El grupo ve el flujo completo de principio a fin.",
        "Actividad 2 — En parejas: mapear un flujo real (50 min). Cada pareja elige un proceso de su área y lo mapea con texto + IA. Identifican 1 cuello de botella y proponen cómo aliviarlo.",
        "Actividad 3 — Individual: documentar un procedimiento propio (30 min) y mini-entregable: 1 SOP + checklist + diagrama de flujo. Que cada alumno comparta el cuello de botella más inesperado que encontró.",
    ],

    "D15": [
        "Arranca M7 sesión 2. Conecta con D14: ya tienen sus procesos documentados. Hoy aplican IA a la cara más visible del trabajo: la atención al cliente. Respuestas coherentes, rápidas y que no dependen de que el experto esté disponible.",
        "El problema de las respuestas variables: cuando cada persona del equipo responde de forma diferente a la misma pregunta, el cliente recibe mensajes contradictorios. Un banco de FAQ resuelve esto. La IA lo construye en minutos.",
        "Privacidad crítica en atención al cliente: NUNCA pegar nombres, DNIs, emails ni datos de contacto reales en ChatGPT. Demostración en vivo: sustituir CLIENTE: María García → CLIENTE_A, EMAIL: maria@empresa.com → EMAIL_OMITIDO. El análisis funciona igual.",
        "FAQ desde tickets ficticios: demuestra el prompt de extracción de preguntas frecuentes desde un lote de consultas anonimizadas. El resultado es un banco de FAQ en tabla con pregunta, respuesta y canal.",
        "El patrón RIA (Reconoce / Informa / Actúa): retoma el patrón de D06, ahora aplicado a atención al cliente. Muestra la diferencia entre una respuesta que ignora el problema, una que informa sin actuar, y una completa con RIA.",
        "Guiones multicanal: el mismo mensaje tiene tono diferente según el canal. Demuestra cómo adaptar la respuesta de un ticket por email, un mensaje de Teams y un guion de llamada telefónica para la misma situación.",
        "Escalado: cuándo NO responder con IA y escalar a un humano. Las 3 situaciones: queja legal, dato no verificable, situación emocional extrema. La IA puede preparar el contexto para el escalado, pero la decisión final es humana.",
        "Actividad 1 — Demo FAQ desde tickets ficticios (formador, 15 min): 10 consultas → banco de 8 FAQ en tabla. El grupo ve el flujo completo.",
        "Actividad 2 — Role-play en parejas (50 min): cliente/agente con IA en vivo. Una persona hace de cliente con queja difícil; la otra usa ChatGPT para preparar la respuesta en tiempo real. Luego cambian roles.",
        "Actividad 3 — Individual: banco de respuestas de su área (30 min) y mini-entregable: banco FAQ + 3 guiones de situaciones difíciles. Que cada pareja comparta la situación de role-play más difícil y cómo la resolvió.",
    ],

    "D16": [
        "Arranca el módulo M8 — Proyecto final. Hoy es el día de construcción individual. Cada alumno tiene una tarea/necesidad real recurrente de su puesto (sembrada en D01, fijada en D09) y hoy la convierte en una solución de IA de extremo a extremo.",
        "Orientación inicial (15 min): recuerda la rúbrica del proyecto. 4 criterios: utilidad real (¿resuelve un problema que tienes esta semana?), dominio de la herramienta (¿funciona sin ayuda?), calidad (¿la salida es profesional?), presentación (¿explicas el impacto?).",
        "El lienzo de proyecto: cada alumno completa los 4 bloques en 10 minutos antes de empezar a construir. PROBLEMA → HERRAMIENTA → PROMPT/FLUJO → RESULTADO ESPERADO. Sin lienzo, el sprint se pierde.",
        "Sprint de construcción — tu rol: mentoría 1:1, rotando cada 8-10 minutos. Identifica quién está bloqueado (problema: no arranca el prompting), quién está en la dirección equivocada (problema: ha elegido una herramienta de pago), y quién va bien (refuerza).",
        "Señales de alarma para intervenir: alguien que sigue en el lienzo después de 20 minutos, alguien que abre Canva cuando su proyecto es de texto, alguien que pega datos reales (nombre + empresa + email) en el prompt.",
        "Ronda de feedback en parejas (mitad del sprint): cada alumno muestra su avance a un compañero durante 5 minutos. El compañero hace 1 pregunta y da 1 sugerencia. No es para juzgar, es para desbloquear.",
        "Indicación de tiempo: a falta de 45 minutos para el cierre, avisa a todos: 'A partir de aquí, lo que no funcione lo simplificamos. El objetivo es tener algo que mostrar mañana, no perfecto sino funcionando.'",
        "Borrador de presentación: los últimos 20 minutos se dedican al borrador. Estructura: 1 min problema → 2 min demo en vivo → 1 min impacto → 30 seg próximos pasos. El alumno que no llega puede usar las últimas slides del brief.",
        "Cierre del sprint: ¿cuántos tienen una solución funcionando? ¿Cuántos tienen al menos la demo lista? El objetivo de mañana es presentar algo real. Recuerda: el foco es el impacto, no la complejidad.",
        "Mini-entregable del día: solución de IA personal funcionando + borrador de presentación para D17. Que cada alumno comparta en 30 segundos cuál es su solución.",
    ],

    "D17": [
        "El día de las presentaciones finales. Recuerda la estructura para el grupo antes de empezar: 5 minutos por persona, 2 minutos de preguntas. El formador toma nota de los puntos fuertes para el feedback colectivo al final.",
        "Criterios de feedback entre pares (escríbelos en la pizarra): 2 cosas que funcionan bien + 1 pregunta sincera + 1 sugerencia de mejora. Que todos los observadores tengan su hoja de feedback lista antes de la primera presentación.",
        "Mientras el alumno presenta: tú observas y tomas notas. No interrumpas. Solo interviene si hay un problema técnico grave. El protagonismo es del alumno. Tu feedback viene en la ronda colectiva.",
        "Entre presentaciones: 30 segundos de transición. Agradece la presentación, anuncia el siguiente nombre, recuerda el tiempo. Mantén el ritmo: con 10 alumnos son 70 minutos de presentaciones.",
        "Feedback del formador: al final de todas las presentaciones, dedica 15 minutos a un feedback colectivo. Señala patrones positivos del grupo ('Todos habéis aplicado la anonimización') y 1-2 áreas de mejora comunes.",
        "Plan personal de 30 días: este es el entregable que los alumnos se llevan al trabajo. 3 acciones concretas para la semana 1, 3 para las semanas 2-3, 1 objetivo de mes. Que cada alumno lo comparta en 30 segundos.",
        "Entrega del cheat-sheet: distribuye (físico o digital) el resumen de herramientas gratuitas. Recuerda que el banco de prompts del curso estará siempre disponible en el notebook.",
        "Reflexión grupal (15 min): 3 preguntas para el grupo. ¿Qué herramienta vas a usar desde mañana? ¿Qué tarea te va a ahorrar más tiempo? ¿Qué le contarías a un compañero que no ha hecho el curso?",
        "Cierre del curso: agradece el esfuerzo del grupo. Recuerda que aprender IA es un proceso continuo: en 6 meses las herramientas habrán evolucionado, pero las habilidades de prompting y pensamiento crítico permanecen.",
        "Foto de grupo y despedida. Si el grupo lo permite, compartir los proyectos finales en un documento compartido para que todos puedan ver las soluciones de sus compañeros.",
    ],
}

# ── FUNCIONES ─────────────────────────────────────────────────────────────────

def detect_day(filename):
    """Extrae el código de día (D05, D06...) del nombre del archivo."""
    import re
    basename = os.path.basename(filename).upper()
    match = re.search(r'D(\d{2})', basename)
    if match:
        return f"D{match.group(1)}"
    return None


def logo_already_added(slide):
    """Comprueba si ya hay un logo AllWomen en la slide (posición top-right)."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            left_in = shape.left / 914400
            width_in = shape.width / 914400
            # Si hay una imagen pequeña en la esquina derecha, es el logo
            if width_in < 3 and left_in > 14:
                return True
    return False


def add_logo(slide):
    """Añade el logo AllWomen si no está ya."""
    if logo_already_added(slide):
        return
    if not os.path.exists(LOGO_PATH):
        print(f"  ⚠️  Logo no encontrado: {LOGO_PATH}")
        return
    slide.shapes.add_picture(LOGO_PATH, LOGO_LEFT, LOGO_TOP, LOGO_WIDTH)


def add_notes(slide, note_text):
    """Añade o reemplaza las speaker notes de una slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = note_text


def process_deck(pptx_path):
    day = detect_day(pptx_path)
    if day not in NOTES:
        print(f"⚠️  {os.path.basename(pptx_path)}: no hay notas para '{day}'. "
              f"Días disponibles: {list(NOTES.keys())}")
        print("   Solo se añadirá el logo.")
        day_notes = []
    else:
        day_notes = NOTES[day]

    print(f"\n📂 Procesando {os.path.basename(pptx_path)} → {day}")
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    print(f"   {len(slides)} slides detectadas")

    for i, slide in enumerate(slides):
        # Logo
        add_logo(slide)
        # Notas
        if i < len(day_notes):
            add_notes(slide, day_notes[i])
            print(f"   Slide {i+1}: logo ✅  notas ✅")
        else:
            print(f"   Slide {i+1}: logo ✅  notas — (sin nota definida para esta slide)")

    prs.save(pptx_path)
    print(f"   💾 Guardado: {pptx_path}")


# ── MAIN ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    for path in sys.argv[1:]:
        if not os.path.exists(path):
            print(f"⚠️  Archivo no encontrado: {path}")
            continue
        if not path.lower().endswith(".pptx"):
            print(f"⚠️  No es un .pptx: {path}")
            continue
        process_deck(path)

    print("\n✅ Procesado completado.")
