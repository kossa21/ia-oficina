#!/usr/bin/env python3
"""
completar_decks.py — IA en la Oficina 01/26
Reescribe las speaker notes de los 17 decks de NotebookLM:
- Lee el OCR de cada slide (inventario JSON) para saber qué muestra.
- Matchea cada slide a la subsección del brief con mayor similitud TF-IDF.
- Genera una nota de formador concreta que referencia el contenido visible
  y los puntos clave del brief.

Guarda copias en Decks_completos/. No toca las imágenes.

Uso:
    python3 completar_decks.py            # todos
    python3 completar_decks.py D05        # uno
"""
import os, re, json, sys, shutil, math
from collections import Counter
from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECKS_DIR = os.path.join(ROOT, "Decks")
BRIEFS_DIR = os.path.join(ROOT, "Briefs")
OUT_DIR = os.path.join(ROOT, "Decks_completos")
OCR_JSON = os.path.join(DECKS_DIR, "_inventario_ocr.json")

MAPPING = {
    "D01": "The_Augmented_Desk.pptx",
    "D02": "Mastering_Structured_AI_Prompts.pptx",
    "D03": "La_Cadena_de_Montaje_IA.pptx",
    "D04": "AI_Efficiency_Operating_System.pptx",
    "D05": "The_AI_Information_Prism.pptx",
    "D06": "Ingeniería_de_Mensajes_con_IA.pptx",
    "D07": "AI_Reporting_Blueprint.pptx",
    "D08": "Smart_AI_Workflows.pptx",
    "D09": "Gestión_de_Equipos_con_IA.pptx",
    "D10": "AI_Data_Analysis_Simplified.pptx",
    "D11": "AI_Research_and_Synthesis.pptx",
    "D12": "AI_Presentation_Design.pptx",
    "D13": "Diseño_Visual_Accesible.pptx",
    "D14": "The_AI_Team_Brain.pptx",
    "D15": "Manual_Táctico_de_Soporte_IA.pptx",
    "D16": "Building_Functional_AI_Solutions.pptx",
    "D17": "Transformación_Digital_con_IA.pptx",
}

STOPWORDS = set("""
a al ante bajo cabe con contra de del desde durante en entre hacia hasta
mediante para por según sin so sobre tras la el los las un una unos unas
y o u e ni que se le lo les nos os te me mi tu su sus es son ser fue era
ha han había habrá hay no si sí más menos como cual cuál qué este esta estos
estas ese esa esos esas aquel aquella aquellos aquellos aquellas también solo sólo
muy muchísimo mucho poco bastante tan tanto bien mejor peor casi ya aún
esto eso aquí ahí allí pero aunque cuando donde porqué porque pues entonces
así así  cada algo nada todo todos toda todas otro otra otros otras
ser estar haber tener hacer ir poder deber querer saber decir ver
hacia mediante puede sigues puedes hagas vas debe debes
así pues luego ahora antes después siempre nunca
""".split())

def tokens(text):
    text = (text or "").lower()
    text = re.sub(r"[áÁ]", "a", text)
    text = re.sub(r"[éÉ]", "e", text)
    text = re.sub(r"[íÍ]", "i", text)
    text = re.sub(r"[óÓ]", "o", text)
    text = re.sub(r"[úÚüÜ]", "u", text)
    text = re.sub(r"[ñÑ]", "n", text)
    text = re.sub(r"[^a-z0-9 ]", " ", text)
    return [t for t in text.split() if len(t) > 3 and t not in STOPWORDS]


# ── PARSEO BRIEF ─────────────────────────────────────────────────────────────

def parse_brief(path):
    with open(path) as f:
        raw = f.read()

    m = re.search(r"^## (Módulo \d+, Sesión [^\n]+)", raw, re.M)
    titulo = m.group(1) if m else ""

    m = re.search(r"## 3\. RESUMEN DE LA SESIÓN\s*(.*?)(?=^---|^## )", raw, re.S | re.M)
    resumen = (m.group(1).strip() if m else "")

    m = re.search(r"## 19\. MENSAJE CLAVE[^\n]*\n+> ?\"?(.*?)\"?\n", raw, re.S)
    mensaje = (m.group(1).strip() if m else "").replace("\n", " ")

    m = re.search(r"## 17\. MINI-ENTREGABLE[^\n]*\n(.*?)(?=^---|^## )", raw, re.S | re.M)
    entregable = m.group(1).strip() if m else ""

    # Actividades (11)
    actividades = []
    sec11 = re.search(r"## 11\. ACTIVIDADES EN AULA(.*?)(?=^## 12)", raw, re.S | re.M)
    if sec11:
        for a in re.finditer(r"### (Actividad \d+ — [^\n]+)\n+\*\*Descripción:\*\* ([^\n]+)", sec11.group(1)):
            actividades.append({"nombre": a.group(1).strip(), "desc": a.group(2).strip()})

    # Sub-secciones del 6 (DESARROLLO TEÓRICO) — fuente principal
    secciones = []
    sec6 = re.search(r"## 6\. DESARROLLO TEÓRICO EN PROSA(.*?)(?=^## 7)", raw, re.S | re.M)
    if sec6:
        parts = re.split(r"^### ", sec6.group(1), flags=re.M)
        for p in parts[1:]:
            lines = p.split("\n", 1)
            sub = lines[0].strip()
            content = lines[1] if len(lines) > 1 else ""
            pc = re.search(r"> \*\*Puntos clave:\*\*\n((?:> ?-[^\n]+\n?)+)", content)
            puntos = []
            if pc:
                for line in pc.group(1).split("\n"):
                    line = re.sub(r"^> ?-\s*", "", line).strip()
                    if line:
                        puntos.append(line)
            content_clean = re.sub(r"> \*\*Puntos clave:\*\*.*?(?=\n###|\Z)", "", content, flags=re.S).strip()
            secciones.append({
                "tipo": "teoria",
                "sub": sub,
                "content": content_clean,
                "puntos": puntos,
            })

    # Casos resueltos (8) como secciones adicionales (para slides de ejemplo)
    sec8 = re.search(r"## 8\. EJEMPLOS RESUELTOS PASO A PASO(.*?)(?=^## 9)", raw, re.S | re.M)
    if sec8:
        parts = re.split(r"^### ", sec8.group(1), flags=re.M)
        for p in parts[1:]:
            lines = p.split("\n", 1)
            sub = lines[0].strip()
            content = lines[1] if len(lines) > 1 else ""
            secciones.append({
                "tipo": "caso",
                "sub": sub,
                "content": content[:600],
                "puntos": [],
            })

    # IDF: corpus = unión de todas las secciones (incluida la teoría)
    docs = [tokens(s["sub"] + " " + s["content"] + " " + " ".join(s["puntos"])) for s in secciones]
    df = Counter()
    for d in docs:
        for t in set(d):
            df[t] += 1
    N = len(docs)
    for i, s in enumerate(secciones):
        toks = docs[i]
        tf = Counter(toks)
        s["tfidf"] = {t: (1 + math.log(c)) * math.log((N + 1) / (df[t] + 1)) for t, c in tf.items()}

    return {
        "titulo": titulo,
        "resumen": resumen,
        "mensaje": mensaje,
        "entregable": entregable,
        "actividades": actividades,
        "secciones": secciones,
    }


# ── MATCHING (TF-IDF) ────────────────────────────────────────────────────────

def match_seccion(ocr_text, secciones, used=None):
    """Devuelve la sección con mayor score TF-IDF respecto al OCR.
    `used` es un Counter de cuántas veces ya se usó cada sección (penalizando repeticiones fuerte)."""
    if not secciones:
        return None
    slide_toks = tokens(ocr_text)
    if not slide_toks:
        return None
    scored = []
    for i, s in enumerate(secciones):
        score = sum(s["tfidf"].get(t, 0) for t in slide_toks)
        if used:
            # Penalización fuerte para forzar diversidad
            score *= (1.0 / (1.0 + used[i] * 1.5))
        scored.append((score, i, s))
    scored.sort(reverse=True, key=lambda x: x[0])
    return scored[0] if scored[0][0] > 0 else (None, None, None)


# ── DETECCIÓN TIPO ───────────────────────────────────────────────────────────

RX_PORTADA = re.compile(r"(tu nuevo asistente|guía visual|IA en la Oficina)", re.I)
RX_EJERCICIO = re.compile(r"\b(ejercicio|práctica|taller|tu turno|momento de|cada (alumno|persona|pareja)|en parejas|grupo|individual)\b", re.I)
RX_CIERRE = re.compile(r"\b(próximo paso|misión|recap|mañana|conclusión|llévate|cierre|para la próxima|hoja de ruta|kit de|resumen visual)\b", re.I)
RX_DATOS = re.compile(r"^\s*\d{2,}\s*%|\b\d{2,}\s*%\s+(de|del)\b", re.I | re.M)
RX_MITOS = re.compile(r"\bmito[s:]?\b", re.I)
# Solo reglas de privacidad/anonimización (no cualquier mención de "regla")
RX_REGLAS = re.compile(r"\b(anonimiz|innegociables?|3 reglas|tres reglas|datos personales|RGPD|reglas? de oro|máquina de anonimización)\b", re.I)
RX_HERRAMIENTAS_PANORAMA = re.compile(r"\b(ecosistema|panorama|kit de|tu arsenal|tu ecosistema)\b", re.I)
RX_TABLA = re.compile(r"\b(comparativa|tabla\s+(de|comparativa)|frente a|vs\.?|matriz)\b", re.I)
RX_PROMPT_DEMO = re.compile(r"^>\s|prompt:|R-C-T-F-R|anatomía", re.I | re.M)
RX_CASO = re.compile(r"\bcaso\s*\d+\b", re.I)

def detect_kind(idx, total, ocr):
    if idx == 0:
        return "portada"
    if idx == total - 1:
        return "cierre"
    if RX_CIERRE.search(ocr) and idx >= total - 3:
        return "cierre"
    if RX_CASO.search(ocr):
        return "caso"
    if RX_EJERCICIO.search(ocr):
        return "ejercicio"
    if RX_MITOS.search(ocr):
        return "mitos"
    if RX_REGLAS.search(ocr):
        return "reglas"
    if RX_HERRAMIENTAS_PANORAMA.search(ocr):
        return "herramientas"
    if RX_DATOS.search(ocr) and len(ocr) < 350:
        return "datos"
    if RX_TABLA.search(ocr):
        return "tabla"
    if RX_PROMPT_DEMO.search(ocr):
        return "prompt"
    return "contenido"


# ── HELPERS ──────────────────────────────────────────────────────────────────

def _clean(text):
    """Limpia residuos OCR: elimina fragmentos iniciales sin sentido, junta espacios."""
    text = re.sub(r"\s+", " ", text or "").strip()
    # Si empieza con un fragmento sin mayúscula inicial Y precedido por puntuación rara, recórtalo
    text = re.sub(r"^[^A-ZÁÉÍÓÚÑa-záéíóúñ¿¡«]+", "", text)
    text = re.sub(r"^[a-záéíóúñ]+ón ", "", text)  # ej "ión a..."
    text = re.sub(r"^[a-z]{1,4} ", "", text)  # ej "ivamente..."
    return text.strip()

def first_sentence(text, max_len=200):
    text = _clean(text)
    if not text:
        return ""
    m = re.match(r"(.{40,%d}?[\.!?])(?:\s|$)" % max_len, text)
    return (m.group(1) if m else text[:max_len]).strip()

def two_sentences(text, max_len=420):
    text = _clean(text)
    sents = re.findall(r".{20,220}?[\.!?]", text)
    out = " ".join(sents[:2]).strip()
    return out[:max_len] if out else first_sentence(text, max_len)

# Títulos genéricos que NotebookLM mete como section dividers
GENERIC_TITLES = re.compile(
    r"^(the analytical blueprint|the prompting blueprint|the augmented desk|"
    r"the ai assembly line|mastering structured ai prompts|"
    r"ai office automation|guía visual|tu nuevo asistente|ia en la oficina)$", re.I
)

def ocr_title(ocr):
    """Devuelve el primer encabezado significativo, saltándose section dividers genéricos."""
    candidates = []
    for line in (ocr or "").split("\n"):
        line = line.strip(" /\\:|—-•·.,()[]\"'“”«»")
        if 4 <= len(line) <= 90 and re.search(r"[A-Za-zÁÉÍÓÚáéíóúÑñ]", line):
            candidates.append(line)
    for c in candidates:
        if not GENERIC_TITLES.match(c.strip()):
            return c
    return candidates[0] if candidates else ""


# ── NOTA POR TIPO ────────────────────────────────────────────────────────────

def note_for(kind, day, idx, total, ocr, seccion_pack, brief):
    """seccion_pack = (score, i, seccion) o (None, None, None)."""
    sec = seccion_pack[2] if seccion_pack and seccion_pack[2] else None
    titulo_slide = ocr_title(ocr)
    onscreen = f"En pantalla: {titulo_slide}. " if titulo_slide else ""

    if kind == "portada":
        return (
            f"Apertura de {day}. Bienvenida y energía: el grupo entra con expectativa, "
            f"baja la temperatura emocional con una pregunta directa al aire — "
            f"'¿qué tarea os ha consumido más tiempo esta semana?'. "
            f"Recoge dos o tres respuestas sin entrar a fondo: te servirán para conectar ejemplos durante el día. "
            f"Encuadra la sesión en una frase: {first_sentence(brief['resumen'])} "
            f"Recuerda al grupo que el material está disponible en NotebookLM (pódcast, vídeo, infografía, fichas, cuestionario) "
            f"y que la regla número uno del curso es no pegar nunca datos personales reales."
        )

    if kind == "datos":
        clave = first_sentence(sec["content"]) if sec else ""
        return (
            f"{onscreen}Cifras del día — déjalas respirar 5 segundos antes de hablar. "
            f"No leas la slide en voz alta: lee solo el porcentaje y deja un silencio. "
            f"{clave} "
            f"Pregunta: '¿Habéis hecho la cuenta en vuestro puesto? Si automatizáis una tarea de 20 minutos que repetís 3 veces por semana, son 50 horas al año.' "
            f"Recoge una o dos reacciones y conecta con el siguiente bloque."
        )

    if kind == "mitos":
        return (
            f"{onscreen}Desmonta los mitos uno a uno, pero brevemente — no más de 30 segundos por mito. "
            f"Insiste en el de las alucinaciones: la IA puede afirmar datos falsos con seguridad total, y por eso la verificación humana es no negociable. "
            f"Sobre 'es demasiado complicado': si sabes escribir un WhatsApp, sabes usar IA. "
            f"Sobre 'va a reemplazar mi puesto': automatiza lo repetitivo, libera lo que requiere tu criterio. "
            f"Sobre 'tengo que pagar': todas las herramientas del curso son gratis con Gmail. "
            f"Pregunta rápida al grupo: '¿Qué mito habéis oído fuera de aquí?'."
        )

    if kind == "reglas":
        # Si la slide habla de datos personales/anonimización → regla de privacidad
        if re.search(r"\b(personales?|anonimiz|RGPD|DNI|salario|expediente)\b", ocr, re.I):
            return (
                f"{onscreen}Esta slide es NO negociable. Léela despacio y mira a la sala mientras lo haces. "
                f"Regla 1: cero datos personales reales en ninguna IA — nombres, DNI, salarios, expedientes. "
                f"Regla 2: anonimiza siempre antes de pegar — sustituye por [CLIENTE_A], [IMPORTE_X], [REF_DOC]. "
                f"Regla 3: el humano firma, la IA redacta — tú revisas y respondes. "
                f"Pregunta al grupo: '¿En qué tarea de vuestro puesto podríais haberos saltado esta regla hoy mismo sin daros cuenta?'. "
                f"Recoge dos respuestas para anclar."
            )
        # Regla genérica de la sesión — usa el contenido del brief
        clave = two_sentences(sec["content"]) if sec else ""
        return (
            f"{onscreen}Esta regla la marcas como NO negociable y la dejas en pizarra durante toda la sesión. "
            f"{clave} "
            f"Pregunta al grupo: '¿Cuándo os ha pasado lo contrario? ¿Cuándo habéis confiado en algo que la IA dijo con seguridad y luego resultó falso?'. "
            f"Recoge una o dos respuestas y conecta con el siguiente bloque."
        )

    if kind == "herramientas":
        return (
            f"{onscreen}Recorrido rápido del ecosistema gratuito: ChatGPT para redactar/analizar texto, "
            f"Gemini si trabajas en Google Workspace (integrado en Gmail/Docs/Sheets), "
            f"Claude para textos largos y análisis preciso, NotebookLM para estudiar el material del día, "
            f"Perplexity para búsquedas con fuentes citadas, Canva para piezas visuales. "
            f"No hay una mejor absoluta: cada una tiene un punto fuerte. "
            f"Pregunta al grupo: '¿Cuál ya usáis aunque sea de forma esporádica?'. "
            f"Aclara: en este curso no se paga nada y todas se abren con la cuenta personal de Gmail."
        )

    if kind == "tabla":
        clave = first_sentence(sec["content"]) if sec else ""
        return (
            f"{onscreen}Tabla comparativa: léela en voz alta columna por columna, no fila por fila — "
            f"así el contraste se ve mejor. {clave} "
            f"Después de leerla pregunta: '¿Cuál enviaríais a un compañero? ¿Por qué?'. "
            f"Si el grupo duda, ayuda con un ejemplo de su contexto antes de pasar a la siguiente slide."
        )

    if kind == "prompt":
        clave = first_sentence(sec["content"]) if sec else ""
        return (
            f"{onscreen}Muestra el prompt completo en pantalla y diséctalo en voz alta señalando cada componente. "
            f"{clave} "
            f"Pídele al grupo que identifique en voz alta dónde está el Rol, el Contexto, la Tarea, el Formato y las Restricciones. "
            f"Si alguno falla, no corrijas tú: pide que otro lo intente. La memoria muscular del R-C-T-F-R se construye así."
        )

    if kind == "caso":
        clave = two_sentences(sec["content"]) if sec else ""
        # Extraer "Caso X" del OCR
        cm = re.search(r"caso\s*(\d+)[^\n]*", ocr, re.I)
        caso_ref = cm.group(0).strip() if cm else "este caso"
        return (
            f"{onscreen}Trabajamos {caso_ref}. Lee la situación en voz alta y deja que el grupo proponga "
            f"cómo lo abordaría antes de mostrar la solución de la slide. "
            f"{clave} "
            f"Después de revelar el prompt, pregunta: '¿Qué cambiaríais para que encaje en vuestro contexto?'. "
            f"Anima a una persona a copiar el prompt y ejecutarlo en directo si tienes tiempo."
        )

    if kind == "ejercicio":
        act = brief['actividades'][0] if brief['actividades'] else None
        nombre = act['nombre'] if act else "Actividad práctica"
        desc = act['desc'] if act else "Cada persona aplica el contenido a un caso de su puesto."
        return (
            f"{onscreen}Momento de práctica: {nombre}. {desc} "
            f"Tu rol durante la práctica: pasea por el aula, desbloquea al que se atasque sin darle la respuesta hecha, "
            f"observa quién no levanta los dedos del teclado para invitarle a probar otra forma. "
            f"Recuérdales la regla de oro antes de empezar: cero datos personales reales — anonimiza con [CORCHETES]. "
            f"Al cerrar, pide a 2-3 voluntarios que muestren su resultado en pantalla."
        )

    if kind == "cierre":
        ent_clean = re.sub(r"\*\*([^*]+)\*\*", r"\1", brief['entregable'] or "")
        ent_short = first_sentence(ent_clean, 240)
        return (
            f"{onscreen}Cierre de la sesión. Recapitula los tres aprendizajes del día en una frase cada uno y "
            f"pide al grupo que reformule el suyo en voz alta — esto fija la memoria. "
            f"Mensaje que se llevan: «{brief['mensaje']}». "
            f"Mini-entregable de hoy: {ent_short} "
            f"Anuncia la siguiente sesión y la misión: traer un caso real (anonimizado) de su puesto para aplicar lo aprendido."
        )

    # contenido genérico, pero ahora con TF-IDF + referencia OCR
    if sec:
        clave = two_sentences(sec["content"])
        puntos_txt = ""
        if sec["puntos"]:
            puntos_txt = " Puntos a martillear: " + "; ".join(sec["puntos"][:3]) + "."
        cierre = " Pregunta al grupo: '¿En qué parte de vuestro flujo de trabajo encaja esto?'. Recoge una respuesta y avanza."
        return f"{onscreen}{clave}{puntos_txt}{cierre}"
    return f"{onscreen}Comenta esta slide en menos de 90 segundos, vincúlala al ejemplo del brief y avanza."


# ── ESCRITURA ────────────────────────────────────────────────────────────────

def set_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.clear()
    tf.paragraphs[0].text = text


# ── PROCESO ──────────────────────────────────────────────────────────────────

def process_day(day, briefs_cache, ocr_cache):
    fname = MAPPING[day]
    src = os.path.join(DECKS_DIR, fname)
    if not os.path.exists(src):
        print(f"⚠️  {day}: no encontrado {fname}")
        return False

    brief_path = os.path.join(BRIEFS_DIR, f"{day}_Brief_NotebookLM.md")
    if not os.path.exists(brief_path):
        print(f"⚠️  {day}: brief no encontrado")
        return False

    if day not in briefs_cache:
        briefs_cache[day] = parse_brief(brief_path)
    brief = briefs_cache[day]
    slides_ocr = ocr_cache[day]["slides"]

    os.makedirs(OUT_DIR, exist_ok=True)
    out_name = f"{day}_{fname.replace('.pptx','')}_completo.pptx"
    dst = os.path.join(OUT_DIR, out_name)
    shutil.copyfile(src, dst)

    prs = Presentation(dst)
    total = len(prs.slides)
    used = Counter()

    for i, slide in enumerate(prs.slides):
        ocr_txt = slides_ocr[i]["ocr"] if i < len(slides_ocr) else ""
        kind = detect_kind(i, total, ocr_txt)
        pack = match_seccion(ocr_txt, brief["secciones"], used)
        if pack and pack[1] is not None:
            used[pack[1]] += 1
        note = note_for(kind, day, i, total, ocr_txt, pack, brief)
        set_notes(slide, note)

    prs.save(dst)
    print(f"✅ {day}: {out_name} ({total} slides)")
    return True


def main():
    if not os.path.exists(OCR_JSON):
        print(f"❌ Falta inventario OCR: {OCR_JSON}. Ejecuta primero ocr_decks.py")
        sys.exit(1)
    with open(OCR_JSON) as f:
        ocr_cache = json.load(f)

    only = sys.argv[1] if len(sys.argv) > 1 else None
    briefs_cache = {}
    days = [only] if only else list(MAPPING.keys())
    ok = 0
    for day in days:
        if day not in MAPPING:
            print(f"⚠️  Día desconocido: {day}")
            continue
        if process_day(day, briefs_cache, ocr_cache):
            ok += 1
    print(f"\n📦 {ok}/{len(days)} decks completados en {OUT_DIR}")


if __name__ == "__main__":
    main()
