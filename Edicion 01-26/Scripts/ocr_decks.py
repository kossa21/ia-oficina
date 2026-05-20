#!/usr/bin/env python3
"""OCR todos los decks mapeados y guarda inventario en JSON."""
import os, json, io, sys
from pptx import Presentation
from PIL import Image
import pytesseract

DECKS_DIR = os.path.join(os.path.dirname(__file__), "..", "Decks")
OUT = os.path.join(os.path.dirname(__file__), "..", "Decks", "_inventario_ocr.json")

MAPPING = {
    "D01": "The_Augmented_Desk.pptx",
    "D02": "Mastering_Structured_AI_Prompts.pptx",
    "D03": "La_Cadena_de_Montaje_IA.pptx",
    "D04": "AI_Efficiency_Operating_System.pptx",
    "D05": "The_AI_Information_Prism.pptx",
    "D06": "Ingeniería_de_Mensajes_con_IA.pptx",
    "D07": "AI_Reporting_Blueprint.pptx",
    "D08": "Smart_AI_Workflows.pptx",
    "D09": "Gestión_de_Equipos_con_IA.pptx",
    "D10": "AI_Data_Analysis_Simplified.pptx",
    "D11": "AI_Research_and_Synthesis.pptx",
    "D12": "AI_Presentation_Design.pptx",
    "D13": "Diseño_Visual_Accesible.pptx",
    "D14": "The_AI_Team_Brain.pptx",
    "D15": "Manual_Táctico_de_Soporte_IA.pptx",
    "D16": "Building_Functional_AI_Solutions.pptx",
    "D17": "Transformación_Digital_con_IA.pptx",
}

def ocr_slide(slide):
    """OCR la imagen principal de la slide."""
    for sh in slide.shapes:
        if sh.shape_type == 13 and sh.width / 914400 > 10:
            img = Image.open(io.BytesIO(sh.image.blob))
            return pytesseract.image_to_string(img, lang='spa').strip()
    return ""

def main():
    inv = {}
    for day, fname in MAPPING.items():
        path = os.path.join(DECKS_DIR, fname)
        if not os.path.exists(path):
            print(f"⚠️  {day}: no encontrado {fname}", file=sys.stderr)
            continue
        prs = Presentation(path)
        slides = []
        for i, s in enumerate(prs.slides):
            txt = ocr_slide(s)
            slides.append({"i": i+1, "ocr": txt})
        inv[day] = {"file": fname, "n_slides": len(slides), "slides": slides}
        print(f"{day}: {len(slides)} slides OCR ✓", file=sys.stderr)
    with open(OUT, "w") as f:
        json.dump(inv, f, ensure_ascii=False, indent=2)
    print(f"\nGuardado: {OUT}")

if __name__ == "__main__":
    main()
